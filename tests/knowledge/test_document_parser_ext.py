from app.core.document_parser import DocumentParser


def test_md_parsed_as_text():
    p = DocumentParser()
    out = p.parse_file("# 标题\n正文内容".encode("utf-8"), "note.md")
    assert out["success"] and out["doc_format"] == "md" and "标题" in out["text"]


def test_pptx_parsed(tmp_path):
    p = DocumentParser()
    try:
        from pptx import Presentation
    except ImportError:
        import pytest
        pytest.skip("python-pptx 未安装")
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "演示要点一"
    path = tmp_path / "s.pptx"
    prs.save(str(path))
    out = p.parse_file(path.read_bytes(), "s.pptx")
    assert out["success"] and out["doc_format"] == "pptx" and "演示要点一" in out["text"]


def test_xlsx_parsed(tmp_path):
    p = DocumentParser()
    try:
        from openpyxl import Workbook
    except ImportError:
        import pytest
        pytest.skip("openpyxl 未安装")
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.append(["姓名", "分数"])
    ws.append(["张三", 90])
    path = tmp_path / "d.xlsx"
    wb.save(str(path))
    out = p.parse_file(path.read_bytes(), "d.xlsx")
    assert out["success"] and out["doc_format"] == "xlsx" and "张三" in out["text"]


def test_missing_dep_returns_failure_not_exception():
    p = DocumentParser()
    try:
        import pptx  # noqa
    except ImportError:
        out = p.parse_file(b"dummy", "x.pptx")
        assert out["success"] is False and "doc_format" in out

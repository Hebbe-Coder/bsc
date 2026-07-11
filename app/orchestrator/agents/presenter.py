# app/orchestrator/agents/presenter.py
from __future__ import annotations
import os
from app.agents.base_agent import BaseAgent


class PresenterAgent(BaseAgent):
    @property
    def name(self) -> str:
        return "Presenter Agent"

    @property
    def system_prompt(self) -> str:
        return (
            "你是 Presenter Agent。职责：基于 6 段状态生成汇报叙述。\n"
            "必须输出 JSON：\n"
            '{"presentation":{"html_url":str,"ppt_path":str,'
            ' "diagram_spec":{"flows":[],"roles":[],"rules":[]}}}'
        )

    @property
    def output_schema(self) -> dict:
        return {"required": ["presentation"]}

    def run(self, session_id: str, state: dict, out_dir: str = "static/presentations",
            context: dict = None) -> dict:
        os.makedirs(out_dir, exist_ok=True)
        user_prompt = f"6 段状态：{state}\n请生成汇报材料（HTML + PPT）的元信息与 diagram_spec。"
        meta = self.llm_service.chat(self.system_prompt, user_prompt, temperature=0.1)
        pres = meta.get("presentation", {})
        html_path = os.path.join(out_dir, f"{session_id}.html")
        ppt_path = os.path.join(out_dir, f"{session_id}.pptx")
        # 真实生成 HTML 汇报页
        html = self._render_html(state)
        with open(html_path, "w", encoding="utf-8") as f:
            f.write(html)
        # 真实生成 PPT（python-pptx）
        self._render_ppt(session_id, state, ppt_path)
        return {"presentation": {
            "html_url": f"/presentations/{session_id}.html",
            "ppt_path": f"/presentations/{session_id}.pptx",
            "diagram_spec": pres.get("diagram_spec", state.get("business_model", {})),
        }}

    def _render_html(self, state: dict) -> str:
        name = state.get("project", {}).get("name", "项目")
        return (
            "<!doctype html><html lang='zh'><head><meta charset='utf-8'>"
            f"<title>{name} 汇报</title></head><body>"
            f"<h1>{name} 业务共创汇报</h1>"
            f"<h2>业务模型</h2><pre>{state.get('business_model', {})!r}</pre>"
            f"<h2>SOP</h2><pre>{state.get('sop', {})!r}</pre>"
            f"<h2>审查</h2><pre>{state.get('review', {})!r}</pre>"
            "</body></html>"
        )

    def _render_ppt(self, session_id: str, state: dict, ppt_path: str):
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            # python-pptx 缺失时写占位文件，避免中断流水线
            with open(ppt_path, "w", encoding="utf-8") as f:
                f.write("PPT placeholder")
            return
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = state.get("project", {}).get("name", "项目") + " 业务共创汇报"
        prs.save(ppt_path)

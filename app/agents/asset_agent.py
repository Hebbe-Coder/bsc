"""Asset Agent — generates PPTX, HTML, and JSON deliverables from workspace."""
from .protocol import BaseAgent, AgentContext, AgentResult, AgentStatus
import logging, os, json, time

logger = logging.getLogger("bsc.studio.asset")

class AssetAgent(BaseAgent):
    name = "asset"
    description = "Generates PPT, HTML, and JSON deliverables from business workspace"
    capabilities = ["generate", "export", "asset"]

    def on_generate(self, ctx: AgentContext, **params) -> dict:
        output_types = params.get("output_types", ["html", "json"])
        assets = []

        output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "output")
        os.makedirs(output_dir, exist_ok=True)

        ts = str(int(time.time()))
        bs = ctx.business_system or {}

        # 1. JSON asset
        if "json" in output_types:
            json_path = os.path.join(output_dir, f"report_{ts}.json")
            with open(json_path, "w", encoding="utf-8") as f:
                json.dump(bs, f, ensure_ascii=False, indent=2)
            assets.append({"type": "json", "path": f"/output/report_{ts}.json", "filename": f"report_{ts}.json"})

        # 2. HTML asset
        if "html" in output_types:
            html_content = self._build_html(bs, ctx.project_name, ctx.domain)
            html_path = os.path.join(output_dir, f"report_{ts}.html")
            with open(html_path, "w", encoding="utf-8") as f:
                f.write(html_content)
            assets.append({"type": "html", "path": f"/output/report_{ts}.html", "filename": f"report_{ts}.html"})

        # 3. PPT placeholder (uses existing generate_ppt.py or mock)
        if "ppt" in output_types:
            ppt_path = os.path.join(output_dir, f"report_{ts}.pptx")
            try:
                self._generate_pptx(bs, ppt_path)
                assets.append({"type": "ppt", "path": f"/output/report_{ts}.pptx", "filename": f"report_{ts}.pptx"})
            except Exception as e:
                logger.warning(f"PPTX generation skipped: {e}")
                assets.append({"type": "ppt", "path": "", "filename": "", "error": str(e)[:100]})

        # Store on context for API response
        if hasattr(ctx, "assets"):
            ctx.assets["assets"] = assets
        else:
            ctx.assets = {"assets": assets}

        return {"assets": assets, "count": len(assets)}

    def _build_html(self, bs: dict, project_name: str, domain: str) -> str:
        """Build a professional consulting-grade HTML report."""
        objectives = bs.get("objectives", [])
        roles = bs.get("roles", [])
        processes = bs.get("processes", [])
        metrics = bs.get("metrics", [])
        risks = bs.get("risks", [])

        obj_items = "".join(f'<li>{o.get("name", str(o)) if isinstance(o, dict) else str(o)}</li>' for o in objectives[:5]) or "<li>Business objectives extracted from PRD</li>"
        role_items = "".join(f'<li>{r.get("name", str(r)) if isinstance(r, dict) else str(r)}</li>' for r in roles[:5]) or "<li>Roles identified from business model</li>"
        proc_items = "".join(f'<tr><td>{i+1}</td><td>{p.get("name", str(p)) if isinstance(p, dict) else str(p)}</td><td><span class=\"badge\">active</span></td></tr>' for i, p in enumerate(processes[:8])) or '<tr><td>1</td><td>Process step</td><td><span class="badge">active</span></td></tr>'

        metric_cards = "".join(f'<div class="kpi-card"><h3>{m.get("name", str(m)) if isinstance(m, dict) else str(m)}</h3><div class="kpi-value">--</div><div class="kpi-target">Target: {m.get("target", "TBD") if isinstance(m, dict) else "TBD"}</div></div>' for m in metrics[:6]) or '<div class="kpi-card"><h3>KPI</h3><div class="kpi-value">--</div></div>'

        risk_rows = "".join(f'<tr><td>{r.get("name", str(r)) if isinstance(r, dict) else str(r)}</td><td><span class="risk-level risk-{r.get("level", "medium") if isinstance(r, dict) else "medium"}">{r.get("level", "medium") if isinstance(r, dict) else "medium"}</span></td><td>{r.get("impact", "--") if isinstance(r, dict) else "--"}</td></tr>' for r in risks[:6]) or '<tr><td>Operational Risk</td><td><span class="risk-level risk-medium">medium</span></td><td>--</td></tr>'

        return f'''<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>BSC Report — {project_name}</title>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{font-family:"Inter","Segoe UI",system-ui,sans-serif;background:#0D1117;color:#E6EDF3;line-height:1.6}}
.container{{max-width:1200px;margin:0 auto;padding:40px 24px}}
header{{text-align:center;padding:60px 0 40px;border-bottom:1px solid #30363D}}
header h1{{font-size:2.4rem;font-weight:700;color:#58A6FF;margin-bottom:8px}}
header .meta{{color:#8B949E;font-size:0.95rem}}
.grid2{{display:grid;grid-template-columns:1fr 1fr;gap:24px;margin:32px 0}}
.card{{background:#161B22;border:1px solid #30363D;border-radius:12px;padding:24px}}
.card h2{{font-size:1.1rem;color:#58A6FF;margin-bottom:16px;text-transform:uppercase;letter-spacing:.5px}}
.card ul,.card ol{{padding-left:20px;color:#C9D1D9}}
.card li{{margin-bottom:6px}}
.kpi-grid{{display:grid;grid-template-columns:repeat(auto-fill,minmax(180px,1fr));gap:16px;margin:32px 0}}
.kpi-card{{background:#161B22;border:1px solid #30363D;border-radius:12px;padding:20px;text-align:center}}
.kpi-card h3{{font-size:0.85rem;color:#8B949E;margin-bottom:8px}}
.kpi-value{{font-size:2rem;font-weight:700;color:#58A6FF}}
.kpi-target{{font-size:0.8rem;color:#8B949E;margin-top:4px}}
table{{width:100%;border-collapse:collapse;margin:16px 0}}
th,td{{padding:12px 16px;text-align:left;border-bottom:1px solid #30363D}}
th{{color:#8B949E;font-size:0.8rem;text-transform:uppercase;letter-spacing:.5px}}
.badge{{background:#23863622;color:#3FB950;padding:2px 10px;border-radius:20px;font-size:0.75rem}}
.risk-level{{padding:2px 10px;border-radius:20px;font-size:0.75rem}}
.risk-high{{background:#DA363322;color:#F85149}}
.risk-medium{{background:#D2992222;color:#D29922}}
.risk-low{{background:#23863622;color:#3FB950}}
footer{{text-align:center;padding:40px 0;color:#484F58;font-size:0.85rem;border-top:1px solid #30363D;margin-top:40px}}
</style>
</head>
<body>
<div class="container">
<header>
<h1>{project_name}</h1>
<div class="meta">BSC Studio · Domain: {domain} · Generated on {time.strftime("%Y-%m-%d %H:%M")}</div>
</header>

<section>
<div class="card"><h2>Business Objectives</h2><ul>{obj_items}</ul></div>
</section>

<div class="grid2">
<div class="card"><h2>Roles & Responsibilities</h2><ul>{role_items}</ul></div>
<div class="card"><h2>Business Processes</h2><table><thead><tr><th>#</th><th>Process</th><th>Status</th></tr></thead><tbody>{proc_items}</tbody></table></div>
</div>

<section>
<h2 style="color:#58A6FF;margin:32px 0 16px">KPI Dashboard</h2>
<div class="kpi-grid">{metric_cards}</div>
</section>

<section>
<h2 style="color:#58A6FF;margin:32px 0 16px">Risk Assessment</h2>
<div class="card">
<table><thead><tr><th>Risk</th><th>Level</th><th>Impact</th></tr></thead><tbody>{risk_rows}</tbody></table>
</div>
</section>

<footer>BSC Studio v5.0 · Business System Compiler · AI-Generated Report</footer>
</div>
</body>
</html>'''

    def _generate_pptx(self, bs: dict, output_path: str):
        """Generate a professional PPTX using python-pptx."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        DARK_BG = RGBColor(0x0D, 0x11, 0x17)
        BLUE = RGBColor(0x58, 0xA6, 0xFF)
        WHITE = RGBColor(0xE6, 0xED, 0xF3)
        GRAY = RGBColor(0x8B, 0x94, 0x9E)

        def add_bg(slide):
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = DARK_BG

        # --- Slide 1: Cover ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
        tf = txBox.text_frame
        tf.paragraphs[0].text = "Business System Report"
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLUE
        p2 = tf.add_paragraph()
        p2.text = f"Domain: {bs.get('domain', 'general')} | BSC Studio v5.0"
        p2.font.size = Pt(18)
        p2.font.color.rgb = GRAY

        # --- Slide 2: Objectives ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
        tf = txBox.text_frame
        tf.paragraphs[0].text = "Business Objectives"
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLUE
        for obj in bs.get("objectives", [])[:8]:
            p = tf.add_paragraph()
            p.text = f"• {obj.get('name', str(obj)) if isinstance(obj, dict) else str(obj)}"
            p.font.size = Pt(16)
            p.font.color.rgb = WHITE
            p.space_after = Pt(8)

        # --- Slide 3: Processes ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
        tf = txBox.text_frame
        tf.paragraphs[0].text = "Business Processes"
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLUE
        for i, proc in enumerate(bs.get("processes", [])[:10]):
            name = proc.get("name", str(proc)) if isinstance(proc, dict) else str(proc)
            p = tf.add_paragraph()
            p.text = f"{i+1}. {name}"
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.space_after = Pt(6)

        # --- Slide 4: KPIs ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
        tf = txBox.text_frame
        tf.paragraphs[0].text = "KPI System"
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLUE
        for m in bs.get("metrics", [])[:8]:
            name = m.get("name", str(m)) if isinstance(m, dict) else str(m)
            target = m.get("target", "--") if isinstance(m, dict) else "--"
            p = tf.add_paragraph()
            p.text = f"• {name} — Target: {target}"
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.space_after = Pt(6)

        # --- Slide 5: Risks ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
        tf = txBox.text_frame
        tf.paragraphs[0].text = "Risk Assessment"
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLUE
        for r in bs.get("risks", [])[:8]:
            name = r.get("name", str(r)) if isinstance(r, dict) else str(r)
            level = r.get("level", "medium") if isinstance(r, dict) else "medium"
            p = tf.add_paragraph()
            p.text = f"• [{level.upper()}] {name}"
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.space_after = Pt(6)

        prs.save(output_path)
        logger.info(f"PPTX saved to {output_path}")

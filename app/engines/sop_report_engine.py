"""
SOP Report Engine - SOP汇报引擎

将SOP流程数据转换为专业、可分享、可导出的汇报文档。

汇报内容包含：
1. 流程概览 - 总步骤数、角色数量、是否含升级机制
2. 详细流程 - 每步的名称、动作、负责人、输入输出、SLA、风险点
3. 角色职责 - 各角色负责的步骤、职责范围
4. SLA汇总 - 各步骤SLA、总耗时预估
5. 风险评估 - 各步骤风险点、缓解措施
6. 流程图 - 可视化流程结构（支持泳道图）
7. 关键成功因素 - 流程成功的关键要素
8. 度量指标 - 效率、质量、成本衡量标准
9. 里程碑规划 - 流程执行关键节点
10. 成本估算 - 人力和时间成本预估
11. 智能摘要 - LLM生成的汇报核心要点
12. AI优化建议 - LLM生成的改进建议

输出格式：
- HTML：专业排版，适合分享和演示
- Markdown：版本控制友好
- PPTX：演示文稿格式
"""
from __future__ import annotations
import uuid
import logging
from typing import Dict, List, Any, Optional
from datetime import datetime

logger = logging.getLogger(__name__)


class SOPReportEngine:
    """SOP汇报引擎"""
    
    def __init__(self):
        self._llm_service = None
    
    def _get_llm_service(self):
        """延迟加载 SOP 专用 LLM 客户端(OpenAI 兼容,多厂商)。"""
        if getattr(self, "_llm_service", None) is None:
            from app.services.sop_llm_client import SOPLLMClient

            self._llm_service = SOPLLMClient()
        return self._llm_service

    def generate_overview(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成流程概览

        包含：
        - 业务领域
        - 核心目标
        - 总步骤数
        - 角色数量
        - 是否含升级机制
        """
        workflow = business_system.get("workflow", [])
        roles = business_system.get("roles", [])
        objectives = business_system.get("objectives", [])
        sla = business_system.get("sla", [])
        
        has_escalation = any(step.get("action", "").lower().find("升级") != -1 for step in workflow)
        
        return {
            "title": "流程概览",
            "description": "SOP流程的整体概况和关键指标",
            "business_domain": business_system.get("business_domain", ""),
            "core_objectives": [obj.get("objective", "") for obj in objectives[:3]],
            "total_steps": len(workflow),
            "total_roles": len(roles),
            "total_sla_items": len(sla),
            "has_escalation": has_escalation,
            "estimated_duration": self._estimate_total_duration(workflow),
        }

    def generate_workflow_detail(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成详细流程

        包含每个步骤的：
        - 步骤编号、名称、动作
        - 负责人/角色
        - 输入/输出
        - SLA
        - 风险点
        """
        workflow = business_system.get("workflow", [])
        risks = business_system.get("risks", [])
        
        steps = []
        for step in workflow:
            step_risks = self._get_step_risks(step, risks)
            
            steps.append({
                "step": step.get("step", ""),
                "name": step.get("name", ""),
                "action": step.get("action", ""),
                "role": step.get("role", "") or step.get("owner", ""),
                "input": step.get("input", ""),
                "output": step.get("output", ""),
                "sla": step.get("sla", ""),
                "risks": [r.get("risk", "") for r in step_risks],
                "mitigations": [r.get("mitigation", "") for r in step_risks],
            })
        
        return {
            "title": "详细流程",
            "description": "每个步骤的具体信息和执行要求",
            "steps": steps,
            "total_steps": len(steps),
        }

    def generate_role_responsibilities(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成角色职责

        包含：
        - 角色名称、部门、级别、人数
        - 负责的步骤
        - 职责描述
        """
        roles = business_system.get("roles", [])
        responsibilities = business_system.get("responsibilities", [])
        workflow = business_system.get("workflow", [])
        
        role_data = []
        
        for role in roles:
            role_name = role.get("role", "")
            
            responsible_steps = []
            for step in workflow:
                if step.get("role") == role_name or step.get("owner") == role_name:
                    responsible_steps.append({
                        "step": step.get("step", ""),
                        "name": step.get("name", ""),
                    })
            
            role_responsibilities = []
            for resp in responsibilities:
                if resp.get("role") == role_name:
                    role_responsibilities = resp.get("duties", [])
                    break
            
            role_data.append({
                "name": role_name,
                "department": role.get("department", ""),
                "level": role.get("level", ""),
                "headcount": role.get("headcount", 1),
                "responsible_steps": responsible_steps,
                "responsibilities": role_responsibilities,
            })
        
        return {
            "title": "角色职责",
            "description": "各角色的职责范围和负责步骤",
            "roles": role_data,
            "total_roles": len(role_data),
        }

    def generate_sla_summary(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成SLA汇总

        包含：
        - 各指标的SLA目标
        - 负责人
        - 总耗时预估
        """
        sla = business_system.get("sla", [])
        kpi = business_system.get("kpi", [])
        workflow = business_system.get("workflow", [])
        
        sla_items = []
        
        for item in sla:
            sla_items.append({
                "metric": item.get("metric", ""),
                "target": item.get("target", ""),
                "owner": item.get("owner", ""),
                "type": "SLA",
            })
        
        for item in kpi:
            sla_items.append({
                "metric": item.get("name", ""),
                "target": item.get("target", ""),
                "owner": item.get("owner", ""),
                "formula": item.get("formula", ""),
                "type": "KPI",
            })
        
        step_slas = []
        for step in workflow:
            if step.get("sla"):
                step_slas.append({
                    "step": step.get("step", ""),
                    "name": step.get("name", ""),
                    "sla": step.get("sla", ""),
                })
        
        return {
            "title": "SLA汇总",
            "description": "服务等级协议和关键绩效指标",
            "sla_items": sla_items,
            "step_slas": step_slas,
            "total_sla_items": len(sla_items),
            "total_step_slas": len(step_slas),
            "estimated_total_duration": self._estimate_total_duration(workflow),
        }

    def generate_risk_assessment(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成风险评估

        包含：
        - 各风险项的描述
        - 严重程度
        - 概率
        - 缓解措施
        """
        risks = business_system.get("risks", [])
        risk_breakdown = business_system.get("risk", {})
        
        all_risks = []
        
        if risks:
            for risk in risks:
                all_risks.append({
                    "risk": risk.get("risk", ""),
                    "severity": risk.get("severity", "") or risk.get("level", ""),
                    "probability": risk.get("probability", ""),
                    "mitigation": risk.get("mitigation", ""),
                    "category": risk.get("category", "其他"),
                })
        else:
            risk_types = [
                ("流程风险", risk_breakdown.get("process_risks", [])),
                ("组织风险", risk_breakdown.get("organization_risks", [])),
                ("系统风险", risk_breakdown.get("system_risks", [])),
                ("合规风险", risk_breakdown.get("compliance_risks", [])),
            ]
            
            for risk_type, items in risk_types:
                for risk in items:
                    all_risks.append({
                        "risk": risk.get("risk", ""),
                        "severity": risk.get("severity", "") or risk.get("level", ""),
                        "probability": risk.get("probability", ""),
                        "mitigation": risk.get("mitigation", ""),
                        "category": risk_type,
                    })
        
        severity_counts = {}
        for risk in all_risks:
            severity = risk["severity"]
            severity_counts[severity] = severity_counts.get(severity, 0) + 1
        
        return {
            "title": "风险评估",
            "description": "流程中的风险点和缓解措施",
            "risks": all_risks,
            "total_risks": len(all_risks),
            "severity_distribution": severity_counts,
        }

    def generate_flowchart(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成流程图数据（支持泳道图）

        包含：
        - 节点列表
        - 连线列表
        - 布局信息
        - 泳道分组信息
        """
        workflow = business_system.get("workflow", [])
        
        nodes = []
        edges = []
        swimlanes = {}
        
        for i, step in enumerate(workflow):
            node_id = f"node_{i+1}"
            role = step.get("role", "") or step.get("owner", "")
            
            nodes.append({
                "id": node_id,
                "step": step.get("step", i+1),
                "name": step.get("name", ""),
                "role": role,
                "type": "process",
            })
            
            if role:
                if role not in swimlanes:
                    swimlanes[role] = []
                swimlanes[role].append(node_id)
            
            if i > 0:
                edges.append({
                    "from": f"node_{i}",
                    "to": node_id,
                    "label": "",
                })
        
        swimlane_data = []
        for role, node_ids in swimlanes.items():
            swimlane_data.append({
                "role": role,
                "nodes": node_ids,
                "node_count": len(node_ids),
            })
        
        return {
            "title": "流程图",
            "description": "可视化的流程结构图（支持泳道图）",
            "nodes": nodes,
            "edges": edges,
            "swimlanes": swimlane_data,
            "total_nodes": len(nodes),
            "total_edges": len(edges),
            "total_swimlanes": len(swimlane_data),
        }

    def generate_gantt_chart(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成甘特图数据

        包含：
        - 任务列表
        - 开始时间
        - 持续时间
        - 依赖关系
        """
        workflow = business_system.get("workflow", [])
        
        tasks = []
        for i, step in enumerate(workflow):
            duration = step.get("sla", "2小时")
            hours = self._parse_duration(duration)
            
            tasks.append({
                "id": f"task_{i+1}",
                "name": step.get("name", f"步骤{i+1}"),
                "step": step.get("step", i+1),
                "role": step.get("role", "") or step.get("owner", ""),
                "duration": hours,
                "duration_text": duration,
                "start_offset": sum(t["duration"] for t in tasks),
                "status": "pending",
            })
        
        return {
            "title": "甘特图",
            "description": "流程执行时间线和任务依赖",
            "tasks": tasks,
            "total_tasks": len(tasks),
            "total_duration": sum(t["duration"] for t in tasks),
        }

    def generate_role_matrix(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成角色职责矩阵

        包含：
        - 角色列表
        - 步骤列表
        - 角色与步骤的对应关系
        """
        workflow = business_system.get("workflow", [])
        roles = business_system.get("roles", [])
        
        role_names = [r.get("role", "") for r in roles]
        step_names = [step.get("name", f"步骤{i+1}") for i, step in enumerate(workflow)]
        
        matrix = []
        for role in role_names:
            row = {"role": role, "responsibilities": []}
            for step in workflow:
                step_role = step.get("role", "") or step.get("owner", "")
                if step_role == role:
                    row["responsibilities"].append({
                        "step": step.get("step", ""),
                        "name": step.get("name", ""),
                        "action": step.get("action", ""),
                        "type": "负责",
                    })
                else:
                    row["responsibilities"].append(None)
            matrix.append(row)
        
        return {
            "title": "角色职责矩阵",
            "description": "角色与步骤的职责对应关系",
            "roles": role_names,
            "steps": step_names,
            "matrix": matrix,
            "total_roles": len(role_names),
            "total_steps": len(step_names),
        }

    def generate_risk_heatmap(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成风险热力图数据

        包含：
        - 风险点列表
        - 概率和严重程度映射
        - 热力图数据
        """
        risks = business_system.get("risks", [])
        risk_breakdown = business_system.get("risk", {})
        
        all_risks = []
        if risks:
            for risk in risks:
                severity = risk.get("severity", "中") or risk.get("level", "中")
                probability = risk.get("probability", "中")
                
                severity_score = {"高": 3, "中": 2, "低": 1}.get(severity, 2)
                probability_score = {"高": 3, "中": 2, "低": 1}.get(probability, 2)
                
                all_risks.append({
                    "risk": risk.get("risk", ""),
                    "severity": severity,
                    "probability": probability,
                    "severity_score": severity_score,
                    "probability_score": probability_score,
                    "risk_score": severity_score * probability_score,
                    "category": risk.get("category", "其他"),
                })
        else:
            risk_types = [
                ("流程风险", risk_breakdown.get("process_risks", [])),
                ("组织风险", risk_breakdown.get("organization_risks", [])),
                ("系统风险", risk_breakdown.get("system_risks", [])),
                ("合规风险", risk_breakdown.get("compliance_risks", [])),
            ]
            
            for risk_type, items in risk_types:
                for risk in items:
                    severity = risk.get("severity", "中") or risk.get("level", "中")
                    probability = risk.get("probability", "中")
                    
                    severity_score = {"高": 3, "中": 2, "低": 1}.get(severity, 2)
                    probability_score = {"高": 3, "中": 2, "低": 1}.get(probability, 2)
                    
                    all_risks.append({
                        "risk": risk.get("risk", ""),
                        "severity": severity,
                        "probability": probability,
                        "severity_score": severity_score,
                        "probability_score": probability_score,
                        "risk_score": severity_score * probability_score,
                        "category": risk_type,
                    })
        
        heatmap_data = [[0, 0, 0], [0, 0, 0], [0, 0, 0]]
        for risk in all_risks:
            heatmap_data[3 - risk["severity_score"]][3 - risk["probability_score"]] += 1
        
        return {
            "title": "风险热力图",
            "description": "风险概率与严重程度的二维分布",
            "risks": all_risks,
            "heatmap": heatmap_data,
            "total_risks": len(all_risks),
            "high_risk_count": sum(1 for r in all_risks if r["risk_score"] >= 6),
        }

    def generate_csf(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成关键成功因素(CSF)

        包含：
        - 关键成功因素列表
        - 影响程度
        - 当前状态
        - 改进行动项
        """
        workflow = business_system.get("workflow", [])
        roles = business_system.get("roles", [])
        objectives = business_system.get("objectives", [])
        
        csf_templates = [
            {
                "id": "csf_001",
                "name": "流程执行效率",
                "description": "各步骤的执行速度和SLA达成率直接影响整体流程效率",
                "impact": "高",
                "status": "部分满足",
                "actions": ["优化瓶颈步骤", "自动化重复操作", "加强监控预警"],
            },
            {
                "id": "csf_002",
                "name": "团队能力匹配",
                "description": "角色职责清晰，人员配置充足，具备所需专业技能",
                "impact": "高",
                "status": "已满足",
                "actions": ["定期技能培训", "建立知识库", "角色备份机制"],
            },
            {
                "id": "csf_003",
                "name": "风险控制能力",
                "description": "识别潜在风险并制定有效的缓解措施",
                "impact": "中",
                "status": "部分满足",
                "actions": ["完善风险评估机制", "制定应急预案", "定期风险复盘"],
            },
            {
                "id": "csf_004",
                "name": "目标对齐",
                "description": "流程执行与业务目标保持一致",
                "impact": "高",
                "status": "已满足",
                "actions": ["定期目标回顾", "KPI跟踪监控", "反馈机制优化"],
            },
            {
                "id": "csf_005",
                "name": "系统支撑能力",
                "description": "IT系统和工具能够支持流程高效执行",
                "impact": "中",
                "status": "部分满足",
                "actions": ["系统性能优化", "工具集成", "用户体验改进"],
            },
        ]
        
        return {
            "title": "关键成功因素",
            "description": "流程成功执行的关键要素和保障措施",
            "factors": csf_templates,
            "total_factors": len(csf_templates),
        }

    def generate_metrics(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成度量指标

        包含：
        - 效率指标
        - 质量指标
        - 成本指标
        """
        workflow = business_system.get("workflow", [])
        sla = business_system.get("sla", [])
        kpi = business_system.get("kpi", [])
        roles = business_system.get("roles", [])
        
        efficiency_metrics = [
            {
                "name": "流程周期时间",
                "current": "N/A",
                "target": self._estimate_total_duration(workflow),
                "unit": "小时",
                "owner": "流程负责人",
            },
            {
                "name": "步骤通过率",
                "current": "N/A",
                "target": ">95%",
                "unit": "%",
                "owner": "各步骤负责人",
            },
            {
                "name": "SLA达成率",
                "current": "N/A",
                "target": ">98%",
                "unit": "%",
                "owner": "流程负责人",
            },
        ]
        
        quality_metrics = [
            {
                "name": "输出准确率",
                "current": "N/A",
                "target": ">99%",
                "unit": "%",
                "owner": "质量保证",
            },
            {
                "name": "客户满意度",
                "current": "N/A",
                "target": ">90%",
                "unit": "%",
                "owner": "客服部门",
            },
            {
                "name": "错误返工率",
                "current": "N/A",
                "target": "<5%",
                "unit": "%",
                "owner": "各步骤负责人",
            },
        ]
        
        cost_metrics = [
            {
                "name": "单位流程成本",
                "current": "N/A",
                "target": "按需优化",
                "unit": "元/次",
                "owner": "财务部门",
            },
            {
                "name": "人力投入",
                "current": f"{len(roles)}人",
                "target": "按需配置",
                "unit": "人",
                "owner": "人力资源",
            },
        ]
        
        return {
            "title": "度量指标",
            "description": "流程效率、质量、成本的衡量标准",
            "efficiency_metrics": efficiency_metrics,
            "quality_metrics": quality_metrics,
            "cost_metrics": cost_metrics,
            "total_efficiency": len(efficiency_metrics),
            "total_quality": len(quality_metrics),
            "total_cost": len(cost_metrics),
        }

    def generate_milestones(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成里程碑规划

        包含：
        - 里程碑节点列表
        - 步骤范围
        - 截止时间
        - 状态
        """
        workflow = business_system.get("workflow", [])
        total_steps = len(workflow)
        
        milestones = []
        if total_steps >= 3:
            step_divisions = [
                (1, int(total_steps * 0.3), "启动阶段", "需求确认与准备"),
                (int(total_steps * 0.3) + 1, int(total_steps * 0.7), "执行阶段", "核心流程执行"),
                (int(total_steps * 0.7) + 1, total_steps, "收尾阶段", "验收与交付"),
            ]
            
            for idx, (start, end, phase_name, milestone_name) in enumerate(step_divisions, 1):
                milestones.append({
                    "id": f"m_{idx:03d}",
                    "name": milestone_name,
                    "phase": phase_name,
                    "step_range": f"{start}-{end}",
                    "deadline": f"T+{idx*3}天",
                    "status": "pending",
                })
        else:
            milestones = [
                {
                    "id": "m_001",
                    "name": "流程启动",
                    "phase": "启动阶段",
                    "step_range": "1",
                    "deadline": "T+1天",
                    "status": "pending",
                },
                {
                    "id": "m_002",
                    "name": "流程完成",
                    "phase": "完成阶段",
                    "step_range": f"1-{total_steps}",
                    "deadline": "T+3天",
                    "status": "pending",
                },
            ]
        
        return {
            "title": "里程碑规划",
            "description": "流程执行的关键节点和时间计划",
            "milestones": milestones,
            "total_milestones": len(milestones),
        }

    def generate_cost_estimate(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成成本估算

        包含：
        - 总工时估算
        - 人力成本估算
        - 成本明细
        """
        workflow = business_system.get("workflow", [])
        roles = business_system.get("roles", [])
        
        avg_hours_per_step = 4
        total_hours = len(workflow) * avg_hours_per_step
        
        avg_hourly_rate = 200
        estimated_cost = total_hours * avg_hourly_rate
        
        breakdown = []
        role_hours = {}
        for step in workflow:
            role = step.get("role", "") or step.get("owner", "未指定")
            role_hours[role] = role_hours.get(role, 0) + avg_hours_per_step
        
        for role, hours in role_hours.items():
            breakdown.append({
                "role": role,
                "hours": hours,
                "cost": f"¥{hours * avg_hourly_rate:,}",
                "steps": [],
            })
        
        return {
            "title": "成本估算",
            "description": "流程执行的人力和时间成本预估",
            "total_hours": total_hours,
            "total_fte": round(total_hours / 176, 2),
            "estimated_cost": f"¥{estimated_cost:,}",
            "average_hourly_rate": f"¥{avg_hourly_rate}",
            "breakdown": breakdown,
            "total_roles": len(breakdown),
        }

    def build_sop_context(self, bs: Dict[str, Any], max_chars: int = 4000) -> str:
        """把 business_system 的真实内容压成紧凑结构化文本,供 LLM 接地。"""
        parts = []
        workflow = bs.get("workflow", [])
        if workflow:
            steps = []
            for i, s in enumerate(workflow[:20], 1):
                owner = s.get("owner", s.get("role", "—"))
                dur = s.get("duration", s.get("estimated_time", "—"))
                steps.append(f"{i}. {s.get('name', s.get('step', '步骤'))} (负责人:{owner}, 时长:{dur})")
            parts.append("【流程步骤】\n" + "\n".join(steps))
        roles = bs.get("roles", [])
        if roles:
            rtxt = "; ".join(
                f"{r.get('role', '角色')}({r.get('department', '—')}, {r.get('headcount', '?')}人)"
                for r in roles[:15]
            )
            parts.append("【角色】" + rtxt)
        sla = bs.get("sla", [])
        if sla:
            parts.append("【SLA】" + "; ".join(
                f"{s.get('step', s.get('name', '环节'))}:{s.get('target', '—')}" for s in sla[:10]
            ))
        kpi = bs.get("kpi", [])
        if kpi:
            parts.append("【KPI】" + "; ".join(
                f"{k.get('name', '指标')}:{k.get('target', '—')}" for k in kpi[:10]
            ))
        risks = bs.get("risks", [])
        if risks:
            rk = [
                f"{r.get('risk', r.get('name', '风险'))}(严重度:{r.get('severity', '—')}, 缓解:{r.get('mitigation', '—')})"
                for r in risks[:10]
            ]
            parts.append("【风险】\n" + "\n".join(rk))
        ctx = "\n\n".join(parts)
        if len(ctx) > max_chars:
            ctx = ctx[:max_chars] + "\n...（已截断）"
        return ctx

    def _fallback_summary(self, bs: Dict[str, Any]) -> Dict[str, Any]:
        workflow = bs.get("workflow", [])
        roles = bs.get("roles", [])
        risks = bs.get("risks", [])
        domain = bs.get("business_domain", "业务")
        return {
            "executive_summary": f"{domain}流程共 {len(workflow)} 个步骤、{len(roles)} 个角色参与,需关注效率与风险控制。",
            "key_findings": [
                f"流程包含 {len(workflow)} 个步骤",
                f"涉及 {len(roles)} 个角色 / 部门",
                f"识别到 {len(risks)} 个风险项",
            ],
            "recommendations": ["定期回顾流程执行情况", "加强风险监控与预警", "对高频步骤考虑自动化"],
            "risk_highlights": [r.get("risk", r.get("name", "风险项")) for r in risks[:2]],
        }

    def _fallback_recommendations(self, bs: Dict[str, Any]) -> Dict[str, Any]:
        workflow = bs.get("workflow", [])
        roles = bs.get("roles", [])
        return {
            "optimization_suggestions": [
                {"id": "opt_001", "title": "流程自动化", "description": "针对重复性步骤引入自动化工具", "priority": "高", "estimated_impact": "节省人力成本", "implementation_steps": ["识别步骤", "选工具", "实施"]},
                {"id": "opt_002", "title": "建立监控机制", "description": "开发流程执行监控看板", "priority": "中", "estimated_impact": "缩短问题响应时间", "implementation_steps": ["定指标", "开发", "上线"]},
            ],
            "prioritized_actions": [
                {"action": f"优先自动化 {len(workflow)} 个步骤中的高频环节", "timeline": "1-2个月"},
                {"action": f"为 {len(roles)} 个角色建立职责看板", "timeline": "2-3个月"},
            ],
        }

    def generate_ai_summary(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成智能摘要（LLM驱动）

        包含：
        - 执行摘要
        - 关键发现
        - 建议
        - 风险亮点
        """
        domain = business_system.get("business_domain", "业务")
        context = self.build_sop_context(business_system)
        system_prompt = (
            "你是资深业务流程分析师。仅基于提供的流程数据,输出严格 JSON,不要任何解释性文字。"
            "字段:executive_summary(一句话核心摘要),key_findings(3-5条字符串),"
            "recommendations(2-3条字符串),risk_highlights(1-2条字符串)。"
        )
        user_prompt = f"业务领域:{domain}\n\n流程数据:\n{context}\n\n请生成执行摘要(JSON)。"
        try:
            client = self._get_llm_service()
            data = client.chat_structured(system_prompt, user_prompt, temperature=0.3, max_tokens=1200)
            if data is None:
                data = self._fallback_summary(business_system)
            return {
                "title": "智能摘要",
                "description": "LLM生成的汇报核心要点",
                "executive_summary": data.get("executive_summary", ""),
                "key_findings": data.get("key_findings", []),
                "recommendations": data.get("recommendations", []),
                "risk_highlights": data.get("risk_highlights", []),
            }
        except Exception as e:
            logger.warning(f"AI摘要生成异常,使用兜底: {e}")
            fb = self._fallback_summary(business_system)
            return {"title": "智能摘要", "description": "LLM生成的汇报核心要点", **fb}

    def generate_ai_recommendations(self, business_system: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成AI优化建议（LLM驱动）

        包含：
        - 优化建议列表
        - 优先级排序的行动项
        """
        domain = business_system.get("business_domain", "业务")
        context = self.build_sop_context(business_system)
        system_prompt = (
            "你是流程优化专家。仅基于提供的流程数据,输出严格 JSON,不要任何解释性文字。"
            "字段:optimization_suggestions(数组,每项 {id,title,description,priority,estimated_impact,implementation_steps[]}),"
            "prioritized_actions(数组,每项 {action,timeline})。"
        )
        user_prompt = f"业务领域:{domain}\n\n流程数据:\n{context}\n\n请提出优化建议(JSON)。"
        try:
            client = self._get_llm_service()
            data = client.chat_structured(system_prompt, user_prompt, temperature=0.5, max_tokens=2000)
            if data is None:
                data = self._fallback_recommendations(business_system)
            return {
                "title": "AI优化建议",
                "description": "LLM生成的流程改进建议",
                "optimization_suggestions": data.get("optimization_suggestions", []),
                "prioritized_actions": data.get("prioritized_actions", []),
            }
        except Exception as e:
            logger.warning(f"AI优化建议生成异常,使用兜底: {e}")
            fb = self._fallback_recommendations(business_system)
            return {"title": "AI优化建议", "description": "LLM生成的流程改进建议", **fb}

    def generate_full_sop_report(self, business_system: Dict[str, Any], enable_ai_analysis: bool = False) -> Dict[str, Any]:
        """
        生成完整的SOP汇报

        包含：
        1. 流程概览
        2. 详细流程
        3. 角色职责
        4. SLA汇总
        5. 风险评估
        6. 流程图数据
        7. 关键成功因素（新增）
        8. 度量指标（新增）
        9. 里程碑规划（新增）
        10. 成本估算（新增）
        11. 智能摘要（新增，需启用AI分析）
        12. AI优化建议（新增，需启用AI分析）

        Args:
            business_system: 业务系统数据
            enable_ai_analysis: 是否启用AI分析功能

        Returns:
            完整的SOP汇报数据
        """
        report = {
            "title": f"{business_system.get('business_domain', '业务系统')}SOP汇报",
            "generated_at": datetime.now().isoformat(),
            "overview": self.generate_overview(business_system),
            "workflow_detail": self.generate_workflow_detail(business_system),
            "role_responsibilities": self.generate_role_responsibilities(business_system),
            "sla_summary": self.generate_sla_summary(business_system),
            "risk_assessment": self.generate_risk_assessment(business_system),
            "flowchart": self.generate_flowchart(business_system),
            "csf": self.generate_csf(business_system),
            "metrics": self.generate_metrics(business_system),
            "milestones": self.generate_milestones(business_system),
            "cost_estimate": self.generate_cost_estimate(business_system),
        }
        
        if enable_ai_analysis:
            report["ai_summary"] = self.generate_ai_summary(business_system)
            report["ai_recommendations"] = self.generate_ai_recommendations(business_system)
        
        return report

    def export_to_markdown(self, report: Dict[str, Any]) -> str:
        """将SOP汇报导出为Markdown格式"""
        md_lines = []
        md_lines.append(f"# {report['title']}")
        md_lines.append(f"生成时间: {report['generated_at']}")
        md_lines.append("")
        
        md_lines.append("---")
        md_lines.append("")
        
        md_lines.append("## 1. 流程概览")
        overview = report["overview"]
        md_lines.append(f"- **业务领域**: {overview['business_domain']}")
        md_lines.append(f"- **总步骤数**: {overview['total_steps']}")
        md_lines.append(f"- **角色数量**: {overview['total_roles']}")
        md_lines.append(f"- **SLA项数**: {overview['total_sla_items']}")
        md_lines.append(f"- **含升级机制**: {'是' if overview['has_escalation'] else '否'}")
        md_lines.append(f"- **预估总耗时**: {overview['estimated_duration']}")
        md_lines.append("")
        
        if overview["core_objectives"]:
            md_lines.append("**核心目标**:")
            for obj in overview["core_objectives"]:
                md_lines.append(f"- {obj}")
            md_lines.append("")
        
        md_lines.append("---")
        md_lines.append("")
        
        md_lines.append("## 2. 详细流程")
        workflow = report["workflow_detail"]
        md_lines.append(f"共 {workflow['total_steps']} 个步骤")
        md_lines.append("")
        
        md_lines.append("| 步骤 | 名称 | 动作 | 负责人 | 输入 | 输出 | SLA |")
        md_lines.append("|------|------|------|--------|------|------|-----|")
        for step in workflow["steps"]:
            md_lines.append(f"| {step['step']} | {step['name']} | {step['action']} | {step['role']} | {step['input']} | {step['output']} | {step['sla']} |")
        md_lines.append("")
        
        for step in workflow["steps"]:
            if step["risks"]:
                md_lines.append(f"### 步骤 {step['step']}: {step['name']} - 风险提示")
                for i, risk in enumerate(step["risks"]):
                    mitigation = step["mitigations"][i] if i < len(step["mitigations"]) else ""
                    md_lines.append(f"- **风险**: {risk}")
                    if mitigation:
                        md_lines.append(f"  - **缓解措施**: {mitigation}")
                md_lines.append("")
        
        md_lines.append("---")
        md_lines.append("")
        
        md_lines.append("## 3. 角色职责")
        roles = report["role_responsibilities"]
        md_lines.append(f"共 {roles['total_roles']} 个角色")
        md_lines.append("")
        
        for role in roles["roles"]:
            md_lines.append(f"### {role['name']}")
            md_lines.append(f"- **部门**: {role['department']}")
            md_lines.append(f"- **级别**: {role['level']}")
            md_lines.append(f"- **人数**: {role['headcount']}")
            md_lines.append("")
            
            if role["responsible_steps"]:
                md_lines.append("**负责步骤**:")
                for step in role["responsible_steps"]:
                    md_lines.append(f"- 步骤 {step['step']}: {step['name']}")
                md_lines.append("")
            
            if role["responsibilities"]:
                md_lines.append("**职责描述**:")
                for duty in role["responsibilities"]:
                    md_lines.append(f"- {duty}")
                md_lines.append("")
        
        md_lines.append("---")
        md_lines.append("")
        
        md_lines.append("## 4. SLA汇总")
        sla = report["sla_summary"]
        md_lines.append(f"共 {sla['total_sla_items']} 项SLA指标")
        md_lines.append("")
        
        md_lines.append("| 指标名称 | 目标 | 负责人 | 类型 |")
        md_lines.append("|----------|------|--------|------|")
        for item in sla["sla_items"]:
            md_lines.append(f"| {item['metric']} | {item['target']} | {item['owner']} | {item['type']} |")
        md_lines.append("")
        
        if sla["step_slas"]:
            md_lines.append("**步骤级SLA**:")
            for item in sla["step_slas"]:
                md_lines.append(f"- 步骤 {item['step']} ({item['name']}): {item['sla']}")
            md_lines.append("")
        
        md_lines.append(f"**预估总耗时**: {sla['estimated_total_duration']}")
        md_lines.append("")
        
        md_lines.append("---")
        md_lines.append("")
        
        md_lines.append("## 5. 风险评估")
        risk = report["risk_assessment"]
        md_lines.append(f"共 {risk['total_risks']} 项风险")
        md_lines.append("")
        
        if risk["severity_distribution"]:
            md_lines.append("**风险等级分布**:")
            for severity, count in risk["severity_distribution"].items():
                md_lines.append(f"- {severity}: {count} 项")
            md_lines.append("")
        
        md_lines.append("| 风险描述 | 严重程度 | 概率 | 缓解措施 | 分类 |")
        md_lines.append("|----------|----------|------|----------|------|")
        for r in risk["risks"]:
            md_lines.append(f"| {r['risk']} | {r['severity']} | {r['probability']} | {r['mitigation']} | {r['category']} |")
        md_lines.append("")
        
        md_lines.append("---")
        md_lines.append("")
        
        md_lines.append("## 6. 流程图")
        flowchart = report["flowchart"]
        md_lines.append(f"共 {flowchart['total_nodes']} 个节点，{flowchart['total_edges']} 条连线")
        md_lines.append("")
        
        md_lines.append("```mermaid")
        md_lines.append("flowchart LR")
        for node in flowchart["nodes"]:
            md_lines.append(f"    {node['id']}[{node['name']}]")
        for edge in flowchart["edges"]:
            md_lines.append(f"    {edge['from']} --> {edge['to']}")
        md_lines.append("```")
        md_lines.append("")
        
        if flowchart.get("swimlanes"):
            md_lines.append("**泳道信息**:")
            for lane in flowchart["swimlanes"]:
                md_lines.append(f"- {lane['role']}: {lane['node_count']} 个步骤")
            md_lines.append("")
        
        md_lines.append("---")
        md_lines.append("")
        
        md_lines.append("## 7. 关键成功因素")
        csf = report.get("csf", {})
        md_lines.append(f"共 {csf.get('total_factors', 0)} 项关键成功因素")
        md_lines.append("")
        
        for factor in csf.get("factors", []):
            md_lines.append(f"### {factor['name']}")
            md_lines.append(f"- **影响程度**: {factor['impact']}")
            md_lines.append(f"- **当前状态**: {factor['status']}")
            md_lines.append(f"- **描述**: {factor['description']}")
            if factor.get("actions"):
                md_lines.append("**改进行动项**:")
                for action in factor["actions"]:
                    md_lines.append(f"  - {action}")
            md_lines.append("")
        
        md_lines.append("---")
        md_lines.append("")
        
        md_lines.append("## 8. 度量指标")
        metrics = report.get("metrics", {})
        md_lines.append("")
        
        md_lines.append("### 效率指标")
        md_lines.append("| 指标名称 | 当前值 | 目标值 | 单位 | 负责人 |")
        md_lines.append("|----------|--------|--------|------|--------|")
        for m in metrics.get("efficiency_metrics", []):
            md_lines.append(f"| {m['name']} | {m['current']} | {m['target']} | {m['unit']} | {m['owner']} |")
        md_lines.append("")
        
        md_lines.append("### 质量指标")
        md_lines.append("| 指标名称 | 当前值 | 目标值 | 单位 | 负责人 |")
        md_lines.append("|----------|--------|--------|------|--------|")
        for m in metrics.get("quality_metrics", []):
            md_lines.append(f"| {m['name']} | {m['current']} | {m['target']} | {m['unit']} | {m['owner']} |")
        md_lines.append("")
        
        md_lines.append("### 成本指标")
        md_lines.append("| 指标名称 | 当前值 | 目标值 | 单位 | 负责人 |")
        md_lines.append("|----------|--------|--------|------|--------|")
        for m in metrics.get("cost_metrics", []):
            md_lines.append(f"| {m['name']} | {m['current']} | {m['target']} | {m['unit']} | {m['owner']} |")
        md_lines.append("")
        
        md_lines.append("---")
        md_lines.append("")
        
        md_lines.append("## 9. 里程碑规划")
        milestones = report.get("milestones", {})
        md_lines.append(f"共 {milestones.get('total_milestones', 0)} 个里程碑")
        md_lines.append("")
        
        md_lines.append("| 里程碑 | 阶段 | 步骤范围 | 截止时间 | 状态 |")
        md_lines.append("|--------|------|----------|----------|------|")
        for m in milestones.get("milestones", []):
            md_lines.append(f"| {m['name']} | {m['phase']} | {m['step_range']} | {m['deadline']} | {m['status']} |")
        md_lines.append("")
        
        md_lines.append("---")
        md_lines.append("")
        
        md_lines.append("## 10. 成本估算")
        cost = report.get("cost_estimate", {})
        md_lines.append(f"- **总工时**: {cost.get('total_hours', 0)} 小时")
        md_lines.append(f"- **人力投入**: {cost.get('total_fte', 0)} FTE")
        md_lines.append(f"- **预估成本**: {cost.get('estimated_cost', '')}")
        md_lines.append(f"- **平均小时费率**: {cost.get('average_hourly_rate', '')}")
        md_lines.append("")
        
        if cost.get("breakdown"):
            md_lines.append("**成本明细**:")
            md_lines.append("| 角色 | 工时(小时) | 成本 |")
            md_lines.append("|------|------------|------|")
            for item in cost["breakdown"]:
                md_lines.append(f"| {item['role']} | {item['hours']} | {item['cost']} |")
            md_lines.append("")
        
        md_lines.append("---")
        md_lines.append("")
        
        if "ai_summary" in report:
            md_lines.append("## 11. 智能摘要")
            ai_summary = report["ai_summary"]
            md_lines.append(f"**执行摘要**: {ai_summary.get('executive_summary', '')}")
            md_lines.append("")
            
            if ai_summary.get("key_findings"):
                md_lines.append("**关键发现**:")
                for finding in ai_summary["key_findings"]:
                    md_lines.append(f"- {finding}")
                md_lines.append("")
            
            if ai_summary.get("recommendations"):
                md_lines.append("**建议**:")
                for rec in ai_summary["recommendations"]:
                    md_lines.append(f"- {rec}")
                md_lines.append("")
            
            if ai_summary.get("risk_highlights"):
                md_lines.append("**风险亮点**:")
                for risk in ai_summary["risk_highlights"]:
                    md_lines.append(f"- {risk}")
                md_lines.append("")
            
            md_lines.append("---")
            md_lines.append("")
        
        if "ai_recommendations" in report:
            md_lines.append("## 12. AI优化建议")
            ai_rec = report["ai_recommendations"]
            md_lines.append("")
            
            for opt in ai_rec.get("optimization_suggestions", []):
                md_lines.append(f"### {opt['title']}")
                md_lines.append(f"- **优先级**: {opt['priority']}")
                md_lines.append(f"- **预估效果**: {opt.get('estimated_impact', '')}")
                md_lines.append(f"- **描述**: {opt['description']}")
                if opt.get("implementation_steps"):
                    md_lines.append("**实施步骤**:")
                    for step in opt["implementation_steps"]:
                        md_lines.append(f"  - {step}")
                md_lines.append("")
            
            if ai_rec.get("prioritized_actions"):
                md_lines.append("**优先级行动项**:")
                for action in ai_rec["prioritized_actions"]:
                    md_lines.append(f"- {action['action']} ({action['timeline']})")
                md_lines.append("")
        
        return "\n".join(md_lines)

    def export_to_pptx(self, report: Dict[str, Any], output_path: Optional[str] = None) -> str:
        """
        将SOP汇报导出为PPTX格式
        
        Args:
            report: SOP汇报数据
            output_path: 输出路径，若为None则自动生成临时路径
        
        Returns:
            输出文件路径
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.enum.shapes import MSO_SHAPE
        except ImportError:
            from exporters.errors import ExportDependencyError
            raise ExportDependencyError(
                fmt="pptx",
                missing_package="python-pptx",
                pip_install="pip install python-pptx>=0.6.20"
            )

        import os
        import tempfile
        
        is_temp = output_path is None
        if is_temp:
            fd, output_path = tempfile.mkstemp(suffix=".pptx")
            os.close(fd)
        else:
            out_dir = os.path.dirname(output_path)
            if out_dir:
                os.makedirs(out_dir, exist_ok=True)

        try:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            DARK_BG = RGBColor(0x1e, 0x29, 0x3b)
            PRIMARY = RGBColor(0x63, 0x66, 0xf1)
            SECONDARY = RGBColor(0x8b, 0x5c, 0xf6)
            WHITE = RGBColor(0xf8, 0xf9, 0xfa)
            GRAY = RGBColor(0x64, 0x74, 0x8b)
            LIGHT_BG = RGBColor(0x0f, 0x17, 0x2a)

            def add_slide(layout_idx=6):
                return prs.slides.add_slide(prs.slide_layouts[layout_idx])

            def add_textbox(slide, left, top, width, height, text, size=12, color=WHITE, bold=False, align=None):
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(size)
                p.font.color.rgb = color
                p.font.bold = bold
                if align:
                    p.alignment = align
                return txBox

            def add_shape(slide, left, top, width, height, fill_color):
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                shape.line.fill.background()
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill_color
                return shape

            def add_table(slide, left, top, width, height, rows, cols, data):
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                for i, row in enumerate(data):
                    for j, cell in enumerate(row):
                        table.cell(i, j).text = str(cell)
                        table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(10)
                        table.cell(i, j).text_frame.paragraphs[0].font.color.rgb = WHITE if i == 0 else GRAY
                        table.cell(i, j).text_frame.paragraphs[0].font.bold = i == 0
                return table

            s = add_slide()
            add_shape(s, 0, 0, prs.slide_width, prs.slide_height, DARK_BG)
            add_shape(s, Inches(0), Inches(2), prs.slide_width, Inches(0.1), PRIMARY)
            add_textbox(s, Inches(1), Inches(1.5), Inches(11), Inches(0.8), report["title"], 36, WHITE, True)
            add_textbox(s, Inches(1), Inches(2.8), Inches(11), Inches(0.3), f"生成时间: {report['generated_at']}", 12, GRAY)
            add_textbox(s, Inches(1), Inches(6.5), Inches(11), Inches(0.3), "SOP Report Engine", 10, SECONDARY)

            overview = report["overview"]
            s = add_slide()
            add_shape(s, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
            add_shape(s, Inches(0), Inches(0.8), prs.slide_width, Inches(0.05), PRIMARY)
            add_textbox(s, Inches(1), Inches(0.3), Inches(10), Inches(0.4), "1. 流程概览", 24, WHITE, True)
            
            metrics = [
                ("业务领域", overview["business_domain"], PRIMARY),
                ("总步骤数", str(overview["total_steps"]), RGBColor(0x3b, 0x82, 0xf6)),
                ("角色数量", str(overview["total_roles"]), RGBColor(0x10, 0xb9, 0x81)),
                ("SLA项数", str(overview["total_sla_items"]), SECONDARY),
                ("升级机制", "是" if overview["has_escalation"] else "否", RGBColor(0xf5, 0x9e, 0x0b)),
                ("预估耗时", overview["estimated_duration"], RGBColor(0x8b, 0x5c, 0xf6)),
            ]
            for i, (label, value, color) in enumerate(metrics):
                row, col = divmod(i, 3)
                cx = Inches(1) + col * Inches(4)
                cy = Inches(1.5) + row * Inches(1.2)
                add_shape(s, cx, cy, Inches(3.5), Inches(1), RGBColor(0x1e, 0x29, 0x3b))
                add_textbox(s, cx + Inches(0.2), cy + Inches(0.1), Inches(3.1), Inches(0.3), label, 10, GRAY)
                add_textbox(s, cx + Inches(0.2), cy + Inches(0.4), Inches(3.1), Inches(0.5), value, 24, color, True)

            if overview["core_objectives"]:
                add_textbox(s, Inches(1), Inches(4), Inches(11), Inches(0.3), "核心目标", 14, SECONDARY, True)
                for i, obj in enumerate(overview["core_objectives"][:5]):
                    add_textbox(s, Inches(1), Inches(4.4) + i * Inches(0.4), Inches(11), Inches(0.3), f"• {obj}", 11, WHITE)

            workflow = report["workflow_detail"]
            s = add_slide()
            add_shape(s, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
            add_shape(s, Inches(0), Inches(0.8), prs.slide_width, Inches(0.05), PRIMARY)
            add_textbox(s, Inches(1), Inches(0.3), Inches(10), Inches(0.4), f"2. 详细流程 ({workflow['total_steps']} 步骤)", 24, WHITE, True)

            if workflow["steps"]:
                rows = min(len(workflow["steps"]) + 1, 9)
                table_data = [["步骤", "名称", "负责人", "SLA"]]
                for step in workflow["steps"][:8]:
                    table_data.append([str(step["step"]), step["name"], step["role"], step["sla"]])
                
                add_table(s, Inches(1), Inches(1.2), Inches(11), Inches(5.5), rows, 4, table_data)

            roles = report["role_responsibilities"]
            s = add_slide()
            add_shape(s, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
            add_shape(s, Inches(0), Inches(0.8), prs.slide_width, Inches(0.05), PRIMARY)
            add_textbox(s, Inches(1), Inches(0.3), Inches(10), Inches(0.4), f"3. 角色职责 ({roles['total_roles']} 角色)", 24, WHITE, True)

            for i, role in enumerate(roles["roles"][:6]):
                row, col = divmod(i, 2)
                cx = Inches(1) + col * Inches(6)
                cy = Inches(1.2) + row * Inches(2.1)
                add_shape(s, cx, cy, Inches(5.5), Inches(1.8), RGBColor(0x1e, 0x29, 0x3b))
                add_shape(s, cx, cy, Inches(5.5), Inches(0.05), SECONDARY)
                add_textbox(s, cx + Inches(0.2), cy + Inches(0.1), Inches(5.1), Inches(0.3), role["name"], 14, WHITE, True)
                add_textbox(s, cx + Inches(0.2), cy + Inches(0.5), Inches(5.1), Inches(0.25), f"{role['department']} | {role['level']} | {role['headcount']}人", 10, GRAY)
                
                if role["responsible_steps"]:
                    steps_str = ", ".join([f"步骤{s['step']}" for s in role["responsible_steps"][:3]])
                    add_textbox(s, cx + Inches(0.2), cy + Inches(0.9), Inches(5.1), Inches(0.25), f"负责: {steps_str}", 10, PRIMARY)

            sla = report["sla_summary"]
            s = add_slide()
            add_shape(s, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
            add_shape(s, Inches(0), Inches(0.8), prs.slide_width, Inches(0.05), PRIMARY)
            add_textbox(s, Inches(1), Inches(0.3), Inches(10), Inches(0.4), f"4. SLA汇总 ({sla['total_sla_items']} 项)", 24, WHITE, True)

            if sla["sla_items"]:
                rows = min(len(sla["sla_items"]) + 1, 9)
                table_data = [["指标", "目标", "负责人", "类型"]]
                for item in sla["sla_items"][:8]:
                    table_data.append([item["metric"], item["target"], item["owner"], item["type"]])
                
                add_table(s, Inches(1), Inches(1.2), Inches(11), Inches(5.5), rows, 4, table_data)

            risk = report["risk_assessment"]
            s = add_slide()
            add_shape(s, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
            add_shape(s, Inches(0), Inches(0.8), prs.slide_width, Inches(0.05), PRIMARY)
            add_textbox(s, Inches(1), Inches(0.3), Inches(10), Inches(0.4), f"5. 风险评估 ({risk['total_risks']} 项)", 24, WHITE, True)

            if risk["risks"]:
                rows = min(len(risk["risks"]) + 1, 8)
                table_data = [["风险描述", "严重程度", "概率", "缓解措施"]]
                for r in risk["risks"][:7]:
                    table_data.append([r["risk"][:30], r["severity"], r["probability"], r["mitigation"][:30]])
                
                add_table(s, Inches(1), Inches(1.2), Inches(11), Inches(5.5), rows, 4, table_data)

            flowchart = report["flowchart"]
            s = add_slide()
            add_shape(s, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
            add_shape(s, Inches(0), Inches(0.8), prs.slide_width, Inches(0.05), PRIMARY)
            add_textbox(s, Inches(1), Inches(0.3), Inches(10), Inches(0.4), f"6. 流程图 ({flowchart['total_nodes']} 节点)", 24, WHITE, True)

            if flowchart["nodes"]:
                steps_per_row = 4
                rows = (len(flowchart["nodes"]) + steps_per_row - 1) // steps_per_row
                
                for i, node in enumerate(flowchart["nodes"]):
                    row_idx, col_idx = divmod(i, steps_per_row)
                    cx = Inches(0.5) + col_idx * Inches(3)
                    cy = Inches(1.5) + row_idx * Inches(1.5)
                    
                    add_shape(s, cx, cy, Inches(2.5), Inches(1), PRIMARY)
                    add_textbox(s, cx + Inches(0.1), cy + Inches(0.1), Inches(2.3), Inches(0.3), f"步骤 {node['step']}", 10, WHITE, True)
                    add_textbox(s, cx + Inches(0.1), cy + Inches(0.4), Inches(2.3), Inches(0.4), node["name"], 10, WHITE)
                    
                    if col_idx < steps_per_row - 1 and i < len(flowchart["nodes"]) - 1:
                        add_textbox(s, cx + Inches(2.6), cy + Inches(0.5), Inches(0.3), Inches(0.3), "→", 16, GRAY)

            s = add_slide()
            add_shape(s, 0, 0, prs.slide_width, prs.slide_height, DARK_BG)
            add_shape(s, Inches(0), Inches(3.5), prs.slide_width, Inches(0.1), PRIMARY)
            add_textbox(s, Inches(1), Inches(3), Inches(11), Inches(0.8), "汇报结束", 36, WHITE, True)
            add_textbox(s, Inches(1), Inches(4), Inches(11), Inches(0.3), report["title"], 14, GRAY)
            add_textbox(s, Inches(1), Inches(6.5), Inches(11), Inches(0.3), "Generated by SOP Report Engine", 10, SECONDARY)

            prs.save(output_path)
            
            if os.path.exists(output_path) and os.path.getsize(output_path) == 0:
                os.remove(output_path)
                raise RuntimeError("生成的PPTX文件为空")
            
            return output_path
            
        except Exception as e:
            if is_temp and os.path.exists(output_path):
                os.remove(output_path)
            raise

    def export_to_html(self, report: Dict[str, Any]) -> str:
        """将SOP汇报导出为HTML格式"""
        html_lines = []
        html_lines.append("""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{ font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif; background: #f8fafc; color: #1e293b; }}
        .container {{ max-width: 1200px; margin: 0 auto; padding: 30px; }}
        .header {{ background: linear-gradient(135deg, #6366f1 0%, #8b5cf6 100%); color: white; padding: 40px; border-radius: 16px; margin-bottom: 30px; }}
        .header h1 {{ font-size: 28px; margin-bottom: 10px; }}
        .header .meta {{ opacity: 0.9; font-size: 14px; }}
        .section {{ background: white; border-radius: 12px; padding: 30px; margin-bottom: 24px; box-shadow: 0 2px 8px rgba(0,0,0,0.06); }}
        .section-title {{ color: #4f46e5; font-size: 20px; margin-bottom: 20px; padding-bottom: 10px; border-bottom: 2px solid #e0e7ff; }}
        .overview-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 20px; }}
        .overview-card {{ background: #f0f9ff; padding: 20px; border-radius: 8px; text-align: center; }}
        .overview-card .label {{ color: #64748b; font-size: 13px; margin-bottom: 8px; }}
        .overview-card .value {{ font-size: 28px; font-weight: bold; color: #0f172a; }}
        .overview-card .value.blue {{ color: #3b82f6; }}
        .overview-card .value.green {{ color: #10b981; }}
        .overview-card .value.purple {{ color: #8b5cf6; }}
        .table {{ width: 100%; border-collapse: collapse; margin-top: 15px; }}
        .table th {{ background: #4f46e5; color: white; padding: 12px; text-align: left; font-weight: 600; font-size: 13px; }}
        .table td {{ border: 1px solid #e2e8f0; padding: 12px; font-size: 13px; }}
        .table tr:nth-child(even) {{ background: #f8fafc; }}
        .table tr:hover {{ background: #f0f9ff; }}
        .role-box {{ background: #fffbeb; border-left: 4px solid #f59e0b; padding: 20px; border-radius: 0 8px 8px 0; margin-bottom: 15px; }}
        .role-box h3 {{ color: #92400e; margin-bottom: 10px; }}
        .role-box .info {{ color: #6b7280; font-size: 13px; margin-bottom: 10px; }}
        .risk-box {{ background: #fef2f2; border-left: 4px solid #ef4444; padding: 15px; border-radius: 0 8px 8px 0; margin-bottom: 10px; }}
        .risk-box .risk-title {{ color: #dc2626; font-weight: 600; }}
        .risk-box .risk-mitigation {{ color: #6b7280; font-size: 13px; margin-top: 5px; }}
        .severity-critical {{ background: #dc2626; color: white; padding: 2px 8px; border-radius: 4px; font-size: 12px; }}
        .severity-high {{ background: #f97316; color: white; padding: 2px 8px; border-radius: 4px; font-size: 12px; }}
        .severity-medium {{ background: #eab308; color: white; padding: 2px 8px; border-radius: 4px; font-size: 12px; }}
        .severity-low {{ background: #22c55e; color: white; padding: 2px 8px; border-radius: 4px; font-size: 12px; }}
        .flowchart-container {{ background: #f8fafc; padding: 20px; border-radius: 8px; overflow-x: auto; }}
        .flowchart-node {{ display: inline-block; background: #6366f1; color: white; padding: 15px 25px; border-radius: 8px; margin: 5px; font-size: 14px; text-align: center; min-width: 120px; }}
        .flowchart-arrow {{ display: inline-block; color: #64748b; font-size: 20px; margin: 0 5px; }}
        .tag {{ display: inline-block; background: #e0e7ff; color: #4338ca; padding: 3px 10px; border-radius: 12px; font-size: 12px; margin-right: 8px; }}
        .objective-list {{ margin-top: 15px; }}
        .objective-item {{ background: #ecfdf5; border: 1px solid #10b981; padding: 10px 15px; border-radius: 6px; margin-bottom: 8px; color: #047857; }}
        .sla-summary {{ margin-top: 20px; padding: 20px; background: #f0fdf4; border-radius: 8px; }}
        .sla-summary .total {{ font-size: 24px; font-weight: bold; color: #16a34a; }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>{title}</h1>
            <div class="meta">生成时间: {generated_at}</div>
        </div>""".format(title=report["title"], generated_at=report["generated_at"]))
        
        overview = report["overview"]
        html_lines.append("""
        <div class="section">
            <div class="section-title">1. 流程概览</div>
            <p style="color: #64748b; margin-bottom: 20px;">{description}</p>
            <div class="overview-grid">
                <div class="overview-card">
                    <div class="label">业务领域</div>
                    <div class="value">{business_domain}</div>
                </div>
                <div class="overview-card">
                    <div class="label">总步骤数</div>
                    <div class="value blue">{total_steps}</div>
                </div>
                <div class="overview-card">
                    <div class="label">角色数量</div>
                    <div class="value green">{total_roles}</div>
                </div>
                <div class="overview-card">
                    <div class="label">SLA项数</div>
                    <div class="value purple">{total_sla_items}</div>
                </div>
            </div>
            <div style="margin-top: 20px; display: flex; gap: 20px;">
                <div class="overview-card" style="flex: 1;">
                    <div class="label">含升级机制</div>
                    <div class="value {escalation_color}">{escalation_text}</div>
                </div>
                <div class="overview-card" style="flex: 1;">
                    <div class="label">预估总耗时</div>
                    <div class="value">{estimated_duration}</div>
                </div>
            </div>""".format(
            description=overview["description"],
            business_domain=overview["business_domain"],
            total_steps=overview["total_steps"],
            total_roles=overview["total_roles"],
            total_sla_items=overview["total_sla_items"],
            escalation_color="green" if overview["has_escalation"] else "gray",
            escalation_text="是" if overview["has_escalation"] else "否",
            estimated_duration=overview["estimated_duration"]
        ))
        
        if overview["core_objectives"]:
            html_lines.append("""
            <div class="objective-list">
                <strong style="color: #059669;">核心目标:</strong>""")
            for obj in overview["core_objectives"]:
                html_lines.append(f"<div class='objective-item'>{obj}</div>")
            html_lines.append("""
            </div>""")
        
        html_lines.append("</div>")
        
        workflow = report["workflow_detail"]
        html_lines.append("""
        <div class="section">
            <div class="section-title">2. 详细流程</div>
            <p style="color: #64748b; margin-bottom: 15px;">共 {total_steps} 个步骤</p>
            <table class="table">
                <thead>
                    <tr>
                        <th>步骤</th>
                        <th>名称</th>
                        <th>动作</th>
                        <th>负责人</th>
                        <th>输入</th>
                        <th>输出</th>
                        <th>SLA</th>
                    </tr>
                </thead>
                <tbody>""".format(total_steps=workflow["total_steps"]))
        
        for step in workflow["steps"]:
            html_lines.append("""
                    <tr>
                        <td><strong>{step}</strong></td>
                        <td>{name}</td>
                        <td>{action}</td>
                        <td>{role}</td>
                        <td>{input}</td>
                        <td>{output}</td>
                        <td><span class="tag">{sla}</span></td>
                    </tr>""".format(
                step=step["step"],
                name=step["name"],
                action=step["action"],
                role=step["role"],
                input=step["input"],
                output=step["output"],
                sla=step["sla"]
            ))
        
        html_lines.append("""
                </tbody>
            </table>""")
        
        for step in workflow["steps"]:
            if step["risks"]:
                html_lines.append("""
            <div style="margin-top: 20px;">
                <h3 style="color: #dc2626; font-size: 16px;">步骤 {step}: {name} - 风险提示</h3>""".format(step=step["step"], name=step["name"]))
                for i, risk in enumerate(step["risks"]):
                    mitigation = step["mitigations"][i] if i < len(step["mitigations"]) else ""
                    html_lines.append("""
                <div class="risk-box">
                    <div class="risk-title">{risk}</div>""".format(risk=risk))
                    if mitigation:
                        html_lines.append("""
                    <div class="risk-mitigation"><strong>缓解措施:</strong> {mitigation}</div>""".format(mitigation=mitigation))
                    html_lines.append("</div>")
                html_lines.append("</div>")
        
        html_lines.append("</div>")
        
        roles = report["role_responsibilities"]
        html_lines.append("""
        <div class="section">
            <div class="section-title">3. 角色职责</div>
            <p style="color: #64748b; margin-bottom: 20px;">共 {total_roles} 个角色</p>""".format(total_roles=roles["total_roles"]))
        
        for role in roles["roles"]:
            html_lines.append("""
            <div class="role-box">
                <h3>{name}</h3>
                <div class="info">部门: {department} | 级别: {level} | 人数: {headcount}</div>""".format(
                name=role["name"],
                department=role["department"],
                level=role["level"],
                headcount=role["headcount"]
            ))
            
            if role["responsible_steps"]:
                html_lines.append("""
                <div style="margin-top: 10px;">
                    <strong style="color: #92400e;">负责步骤:</strong>
                    <ul style="margin-left: 20px; margin-top: 5px; color: #4b5563;">""")
                for step in role["responsible_steps"]:
                    html_lines.append(f"<li>步骤 {step['step']}: {step['name']}</li>")
                html_lines.append("</ul></div>")
            
            if role["responsibilities"]:
                html_lines.append("""
                <div style="margin-top: 10px;">
                    <strong style="color: #92400e;">职责描述:</strong>
                    <ul style="margin-left: 20px; margin-top: 5px; color: #4b5563;">""")
                for duty in role["responsibilities"]:
                    html_lines.append(f"<li>{duty}</li>")
                html_lines.append("</ul></div>")
            
            html_lines.append("</div>")
        
        html_lines.append("</div>")
        
        sla = report["sla_summary"]
        html_lines.append("""
        <div class="section">
            <div class="section-title">4. SLA汇总</div>
            <p style="color: #64748b; margin-bottom: 15px;">共 {total_sla_items} 项SLA指标</p>
            <table class="table">
                <thead>
                    <tr>
                        <th>指标名称</th>
                        <th>目标</th>
                        <th>负责人</th>
                        <th>类型</th>
                    </tr>
                </thead>
                <tbody>""".format(total_sla_items=sla["total_sla_items"]))
        
        for item in sla["sla_items"]:
            html_lines.append("""
                    <tr>
                        <td>{metric}</td>
                        <td><strong>{target}</strong></td>
                        <td>{owner}</td>
                        <td><span class="tag">{type}</span></td>
                    </tr>""".format(
                metric=item["metric"],
                target=item["target"],
                owner=item["owner"],
                type=item["type"]
            ))
        
        html_lines.append("""
                </tbody>
            </table>""")
        
        if sla["step_slas"]:
            html_lines.append("""
            <div style="margin-top: 20px;">
                <strong style="color: #16a34a;">步骤级SLA:</strong>
                <div style="margin-top: 10px;">""")
            for item in sla["step_slas"]:
                html_lines.append(f"<div style='margin-bottom: 5px;'>步骤 {item['step']} ({item['name']}): <strong>{item['sla']}</strong></div>")
            html_lines.append("</div></div>")
        
        html_lines.append("""
            <div class="sla-summary">
                <div style="color: #059669;">预估总耗时</div>
                <div class="total">{estimated_total_duration}</div>
            </div>
        </div>""".format(estimated_total_duration=sla["estimated_total_duration"]))
        
        risk = report["risk_assessment"]
        html_lines.append("""
        <div class="section">
            <div class="section-title">5. 风险评估</div>
            <p style="color: #64748b; margin-bottom: 15px;">共 {total_risks} 项风险</p>""".format(total_risks=risk["total_risks"]))
        
        if risk["severity_distribution"]:
            html_lines.append("""
            <div style="margin-bottom: 20px;">
                <strong style="color: #dc2626;">风险等级分布:</strong>
                <div style="margin-top: 10px;">""")
            for severity, count in risk["severity_distribution"].items():
                severity_class = f"severity-{severity.lower()}" if severity.lower() in ["critical", "high", "medium", "low"] else "severity-medium"
                html_lines.append(f"<span class='{severity_class}' style='margin-right: 15px;'>{severity}: {count} 项</span>")
            html_lines.append("</div></div>")
        
        html_lines.append("""
            <table class="table">
                <thead>
                    <tr>
                        <th>风险描述</th>
                        <th>严重程度</th>
                        <th>概率</th>
                        <th>缓解措施</th>
                        <th>分类</th>
                    </tr>
                </thead>
                <tbody>""")
        
        for r in risk["risks"]:
            severity_class = f"severity-{r['severity'].lower()}" if r['severity'].lower() in ["critical", "high", "medium", "low"] else "severity-medium"
            html_lines.append("""
                    <tr>
                        <td>{risk}</td>
                        <td><span class="{severity_class}">{severity}</span></td>
                        <td>{probability}</td>
                        <td>{mitigation}</td>
                        <td><span class="tag">{category}</span></td>
                    </tr>""".format(
                risk=r["risk"],
                severity_class=severity_class,
                severity=r["severity"],
                probability=r["probability"],
                mitigation=r["mitigation"],
                category=r["category"]
            ))
        
        html_lines.append("""
                </tbody>
            </table>
        </div>""")
        
        flowchart = report["flowchart"]
        html_lines.append("""
        <div class="section">
            <div class="section-title">6. 流程图</div>
            <p style="color: #64748b; margin-bottom: 20px;">共 {total_nodes} 个节点，{total_edges} 条连线</p>
            <div class="flowchart-container">""".format(total_nodes=flowchart["total_nodes"], total_edges=flowchart["total_edges"]))
        
        for i, node in enumerate(flowchart["nodes"]):
            html_lines.append(f"<div class='flowchart-node'>步骤 {node['step']}<br>{node['name']}</div>")
            if i < len(flowchart["nodes"]) - 1:
                html_lines.append("<div class='flowchart-arrow'>→</div>")
        
        html_lines.append("""
            </div>
        </div>""")
        
        csf = report.get("csf", {})
        html_lines.append("""
        <div class="section">
            <div class="section-title">7. 关键成功因素</div>
            <p style="color: #64748b; margin-bottom: 20px;">共 {total_factors} 项关键成功因素</p>""".format(total_factors=csf.get("total_factors", 0)))
        
        for factor in csf.get("factors", []):
            impact_color = {"高": "#dc2626", "中": "#f59e0b", "低": "#10b981"}.get(factor["impact"], "#64748b")
            status_color = {"已满足": "#10b981", "部分满足": "#f59e0b", "未满足": "#dc2626"}.get(factor["status"], "#64748b")
            html_lines.append("""
            <div style="background: #f0fdf4; border-left: 4px solid {impact_color}; padding: 20px; border-radius: 0 8px 8px 0; margin-bottom: 15px;">
                <h3 style="color: #166534;">{name}</h3>
                <div style="margin-top: 10px;">
                    <span style="display: inline-block; background: {impact_color}; color: white; padding: 3px 8px; border-radius: 4px; font-size: 12px; margin-right: 10px;">影响程度: {impact}</span>
                    <span style="display: inline-block; background: {status_color}; color: white; padding: 3px 8px; border-radius: 4px; font-size: 12px;">当前状态: {status}</span>
                </div>
                <p style="margin-top: 10px; color: #4b5563;">{description}</p>""".format(
                name=factor["name"],
                impact=factor["impact"],
                status=factor["status"],
                description=factor["description"],
                impact_color=impact_color,
                status_color=status_color
            ))
            if factor.get("actions"):
                html_lines.append("""
                <div style="margin-top: 10px;">
                    <strong style="color: #166534;">改进行动项:</strong>
                    <ul style="margin-left: 20px; margin-top: 5px; color: #4b5563;">""")
                for action in factor["actions"]:
                    html_lines.append(f"<li>{action}</li>")
                html_lines.append("</ul></div>")
            html_lines.append("</div>")
        
        html_lines.append("</div>")
        
        metrics = report.get("metrics", {})
        html_lines.append("""
        <div class="section">
            <div class="section-title">8. 度量指标</div>""")
        
        html_lines.append("""
            <h4 style="color: #3b82f6; margin-top: 20px; margin-bottom: 15px;">效率指标</h4>
            <table class="table">
                <thead>
                    <tr><th>指标名称</th><th>当前值</th><th>目标值</th><th>单位</th><th>负责人</th></tr>
                </thead>
                <tbody>""")
        for m in metrics.get("efficiency_metrics", []):
            html_lines.append(f"""
                    <tr><td>{m['name']}</td><td>{m['current']}</td><td>{m['target']}</td><td>{m['unit']}</td><td>{m['owner']}</td></tr>""")
        html_lines.append("""
                </tbody>
            </table>""")
        
        html_lines.append("""
            <h4 style="color: #10b981; margin-top: 20px; margin-bottom: 15px;">质量指标</h4>
            <table class="table">
                <thead>
                    <tr><th>指标名称</th><th>当前值</th><th>目标值</th><th>单位</th><th>负责人</th></tr>
                </thead>
                <tbody>""")
        for m in metrics.get("quality_metrics", []):
            html_lines.append(f"""
                    <tr><td>{m['name']}</td><td>{m['current']}</td><td>{m['target']}</td><td>{m['unit']}</td><td>{m['owner']}</td></tr>""")
        html_lines.append("""
                </tbody>
            </table>""")
        
        html_lines.append("""
            <h4 style="color: #f59e0b; margin-top: 20px; margin-bottom: 15px;">成本指标</h4>
            <table class="table">
                <thead>
                    <tr><th>指标名称</th><th>当前值</th><th>目标值</th><th>单位</th><th>负责人</th></tr>
                </thead>
                <tbody>""")
        for m in metrics.get("cost_metrics", []):
            html_lines.append(f"""
                    <tr><td>{m['name']}</td><td>{m['current']}</td><td>{m['target']}</td><td>{m['unit']}</td><td>{m['owner']}</td></tr>""")
        html_lines.append("""
                </tbody>
            </table>""")
        
        html_lines.append("</div>")
        
        milestones = report.get("milestones", {})
        html_lines.append("""
        <div class="section">
            <div class="section-title">9. 里程碑规划</div>
            <p style="color: #64748b; margin-bottom: 20px;">共 {total_milestones} 个里程碑</p>
            <table class="table">
                <thead>
                    <tr><th>里程碑</th><th>阶段</th><th>步骤范围</th><th>截止时间</th><th>状态</th></tr>
                </thead>
                <tbody>""".format(total_milestones=milestones.get("total_milestones", 0)))
        for m in milestones.get("milestones", []):
            status_color = {"pending": "#f59e0b", "completed": "#10b981", "blocked": "#dc2626"}.get(m["status"], "#64748b")
            html_lines.append(f"""
                    <tr><td>{m['name']}</td><td>{m['phase']}</td><td>{m['step_range']}</td><td>{m['deadline']}</td><td><span style="color: {status_color}; font-weight: bold;">{m['status']}</span></td></tr>""")
        html_lines.append("""
                </tbody>
            </table>
        </div>""")
        
        cost = report.get("cost_estimate", {})
        html_lines.append("""
        <div class="section">
            <div class="section-title">10. 成本估算</div>
            <div class="overview-grid">
                <div class="overview-card">
                    <div class="label">总工时</div>
                    <div class="value blue">{total_hours} 小时</div>
                </div>
                <div class="overview-card">
                    <div class="label">人力投入</div>
                    <div class="value green">{total_fte} FTE</div>
                </div>
                <div class="overview-card">
                    <div class="label">预估成本</div>
                    <div class="value purple">{estimated_cost}</div>
                </div>
                <div class="overview-card">
                    <div class="label">平均小时费率</div>
                    <div class="value">{average_hourly_rate}</div>
                </div>
            </div>""".format(
            total_hours=cost.get("total_hours", 0),
            total_fte=cost.get("total_fte", 0),
            estimated_cost=cost.get("estimated_cost", ""),
            average_hourly_rate=cost.get("average_hourly_rate", "")
        ))
        
        if cost.get("breakdown"):
            html_lines.append("""
            <h4 style="color: #4f46e5; margin-top: 20px; margin-bottom: 15px;">成本明细</h4>
            <table class="table">
                <thead>
                    <tr><th>角色</th><th>工时(小时)</th><th>成本</th></tr>
                </thead>
                <tbody>""")
            for item in cost["breakdown"]:
                html_lines.append(f"""
                    <tr><td>{item['role']}</td><td>{item['hours']}</td><td>{item['cost']}</td></tr>""")
            html_lines.append("""
                </tbody>
            </table>""")
        
        html_lines.append("</div>")
        
        if "ai_summary" in report:
            ai_summary = report["ai_summary"]
            html_lines.append("""
            <div class="section">
                <div class="section-title">11. 智能摘要</div>
                <div style="background: linear-gradient(135deg, #f0f9ff 0%, #e0f2fe 100%); padding: 20px; border-radius: 8px; margin-bottom: 20px;">
                    <h4 style="color: #0369a1; margin-bottom: 10px;">执行摘要</h4>
                    <p style="color: #075985; font-size: 15px;">{executive_summary}</p>
                </div>""".format(executive_summary=ai_summary.get("executive_summary", "")))
            
            if ai_summary.get("key_findings"):
                html_lines.append("""
                <h4 style="color: #4f46e5; margin-bottom: 15px;">关键发现</h4>
                <ul style="margin-left: 20px; color: #4b5563;">""")
                for finding in ai_summary["key_findings"]:
                    html_lines.append(f"<li>{finding}</li>")
                html_lines.append("</ul>")
            
            if ai_summary.get("recommendations"):
                html_lines.append("""
                <h4 style="color: #10b981; margin-bottom: 15px;">建议</h4>
                <ul style="margin-left: 20px; color: #4b5563;">""")
                for rec in ai_summary["recommendations"]:
                    html_lines.append(f"<li>{rec}</li>")
                html_lines.append("</ul>")
            
            if ai_summary.get("risk_highlights"):
                html_lines.append("""
                <h4 style="color: #dc2626; margin-bottom: 15px;">风险亮点</h4>
                <ul style="margin-left: 20px; color: #4b5563;">""")
                for risk in ai_summary["risk_highlights"]:
                    html_lines.append(f"<li>{risk}</li>")
                html_lines.append("</ul>")
            
            html_lines.append("</div>")
        
        if "ai_recommendations" in report:
            ai_rec = report["ai_recommendations"]
            html_lines.append("""
            <div class="section">
                <div class="section-title">12. AI优化建议</div>""")
            
            for opt in ai_rec.get("optimization_suggestions", []):
                priority_color = {"高": "#dc2626", "中": "#f59e0b", "低": "#10b981"}.get(opt["priority"], "#64748b")
                html_lines.append("""
                <div style="background: #faf5ff; border-left: 4px solid {priority_color}; padding: 20px; border-radius: 0 8px 8px 0; margin-bottom: 15px;">
                    <h3 style="color: #7c3aed;">{title}</h3>
                    <div style="margin-top: 10px;">
                        <span style="display: inline-block; background: {priority_color}; color: white; padding: 3px 8px; border-radius: 4px; font-size: 12px; margin-right: 10px;">优先级: {priority}</span>
                        <span style="color: #64748b;">{estimated_impact}</span>
                    </div>
                    <p style="margin-top: 10px; color: #4b5563;">{description}</p>""".format(
                title=opt["title"],
                priority=opt["priority"],
                estimated_impact=opt.get("estimated_impact", ""),
                description=opt["description"],
                priority_color=priority_color
            ))
                if opt.get("implementation_steps"):
                    html_lines.append("""
                    <div style="margin-top: 10px;">
                        <strong style="color: #7c3aed;">实施步骤:</strong>
                        <ol style="margin-left: 20px; margin-top: 5px; color: #4b5563;">""")
                    for step in opt["implementation_steps"]:
                        html_lines.append(f"<li>{step}</li>")
                    html_lines.append("</ol></div>")
                html_lines.append("</div>")
            
            if ai_rec.get("prioritized_actions"):
                html_lines.append("""
                <h4 style="color: #4f46e5; margin-top: 20px; margin-bottom: 15px;">优先级行动项</h4>
                <ul style="margin-left: 20px; color: #4b5563;">""")
                for action in ai_rec["prioritized_actions"]:
                    html_lines.append(f"<li><strong>{action['action']}</strong> - {action['timeline']}</li>")
                html_lines.append("</ul>")
            
            html_lines.append("</div>")
        
        html_lines.append("""
    </div>
</body>
</html>""")
        
        return "\n".join(html_lines)

    def _estimate_total_duration(self, workflow: List[Dict[str, Any]]) -> str:
        """估算总耗时"""
        duration_mapping = {
            "分钟": 1,
            "小时": 60,
            "天": 1440,
            "周": 10080,
        }
        
        total_minutes = 0
        
        for step in workflow:
            sla = step.get("sla", "")
            for unit, multiplier in duration_mapping.items():
                if unit in sla:
                    import re
                    match = re.search(r'(\d+)\s*' + re.escape(unit), sla)
                    if match:
                        total_minutes += int(match.group(1)) * multiplier
        
        if total_minutes < 60:
            return f"{total_minutes}分钟"
        elif total_minutes < 1440:
            hours = total_minutes // 60
            minutes = total_minutes % 60
            return f"{hours}小时{minutes}分钟" if minutes else f"{hours}小时"
        elif total_minutes < 10080:
            days = total_minutes // 1440
            hours = (total_minutes % 1440) // 60
            return f"{days}天{hours}小时" if hours else f"{days}天"
        else:
            weeks = total_minutes // 10080
            days = (total_minutes % 10080) // 1440
            return f"{weeks}周{days}天" if days else f"{weeks}周"

    def _parse_duration(self, duration_text: str) -> float:
        """解析时长文本为小时数"""
        import re
        
        duration_mapping = {
            "分钟": 1 / 60,
            "小时": 1,
            "天": 24,
            "周": 168,
        }
        
        for unit, multiplier in duration_mapping.items():
            if unit in duration_text:
                match = re.search(r'(\d+\.?\d*)\s*' + re.escape(unit), duration_text)
                if match:
                    return float(match.group(1)) * multiplier
        
        return 2.0

    def _get_step_risks(self, step: Dict[str, Any], risks: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """获取步骤相关的风险"""
        step_name = step.get("name", "")
        step_action = step.get("action", "")
        step_role = step.get("role", "")
        
        step_risks = []
        for risk in risks:
            risk_desc = risk.get("risk", "")
            if step_name in risk_desc or step_action[:10] in risk_desc or step_role in risk_desc:
                step_risks.append(risk)
        
        return step_risks


__all__ = ["SOPReportEngine"]
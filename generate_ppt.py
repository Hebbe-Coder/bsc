"""BSC Studio Investor Pitch - CATARC Modern Design Language"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DN = RGBColor(0x00,0x15,0x29)  # deep navy
TB = RGBColor(0x18,0x90,0xFF)  # tech blue
NC = RGBColor(0x00,0xE5,0xFF)  # neon cyan
PG = RGBColor(0xF7,0xF9,0xFC)  # polar gray
W  = RGBColor(0xFF,0xFF,0xFF)
DT = RGBColor(0x37,0x41,0x51)  # dark text
S  = RGBColor(0x6B,0x72,0x80)  # secondary
G  = RGBColor(0x10,0xB9,0x81)  # green
R  = RGBColor(0xEF,0x44,0x44)  # red
O  = RGBColor(0xF5,0x9E,0x0B)  # orange

SW = Inches(13.333); SH = Inches(7.5)
prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH

def blank(): return prs.slides.add_slide(prs.slide_layouts[6])

def r(slide, l, t, w, h, c, ln=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if c is not None: s.fill.solid(); s.fill.fore_color.rgb = c
    s.line.fill.background()
    if ln is not None: s.line.color.rgb = ln; s.line.width = Pt(1)
    return s

def tx(slide, l, t, w, h, txt, fs=16, c=W, b=False, ff="Microsoft YaHei", a=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(fs); p.font.color.rgb = c; p.font.bold = b
    p.font.name = ff; p.alignment = a
    return tb

def pn(slide, n, total=12):
    tx(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.3), f"{n}/{total}", 10, S, False, "Arial", PP_ALIGN.RIGHT)

def hdr(slide, title, ch="01"):
    r(slide, Inches(0), Inches(0), SW, Inches(0.9), DN)
    r(slide, Inches(0), Inches(0.9), SW, Inches(0.05), TB)
    r(slide, Inches(0), Inches(0), Inches(0.6), Inches(0.9), RGBColor(0x0A,0x0A,0x1A))
    tx(slide, Inches(0.05), Inches(0.2), Inches(0.5), Inches(0.5), ch, 22, NC, True, "Arial", PP_ALIGN.CENTER)
    tx(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.5), title, 24, W, True)
    tx(slide, Inches(10.5), Inches(0.2), Inches(2.5), Inches(0.5), "BSC Studio", 18, TB, True, "Arial", PP_ALIGN.RIGHT)

def kpi(slide, x, y, bw, bh, lab, val, clr=TB):
    ix,iy,iw = int(x),int(y),int(bw)
    r(slide, x, y, bw, bh, W)
    r(slide, x, y, bw, Inches(0.04), clr)
    tx(slide, ix+182880, iy+137160, iw-365760, 228600, lab, 10, S, False, "Arial")
    tx(slide, ix+182880, iy+365760, iw-365760, 457200, str(val), 28, clr, True, "Arial")

# ==== S1: COVER ====
s = blank(); r(s, Inches(0), Inches(0), SW, SH, DN)
r(s, Inches(7), Inches(0), Inches(6.3), SH, RGBColor(0x00,0x08,0x18))
r(s, Inches(1), Inches(1.5), Inches(0.8), Inches(0.06), NC)
tx(s, Inches(1), Inches(2.5), Inches(8), Inches(1.2), "BSC Studio", 64, W, True)
tx(s, Inches(1), Inches(3.5), Inches(8), Inches(0.6), "Business System Compiler", 28, NC, False, "Arial")
tx(s, Inches(1), Inches(4.3), Inches(7), Inches(0.5), "AI-Powered Enterprise Intelligence Platform", 18, RGBColor(0xA0,0xAE,0xC0))
r(s, Inches(1), Inches(6.0), Inches(0.06), Inches(0.8), NC)
tx(s, Inches(1.3), Inches(6.1), Inches(3), Inches(0.3), "Investor Pitch", 16, RGBColor(0xA0,0xAE,0xC0))
tx(s, Inches(1.3), Inches(6.4), Inches(3), Inches(0.3), "July 2026", 20, W, True, "Arial")

# ==== S2: PROBLEM ====
s = blank(); r(s, Inches(0), Inches(0), SW, SH, PG); hdr(s, "The Problem", "01")
tx(s, Inches(1), Inches(2.0), Inches(5), Inches(1.5), "Business documents are silos of unstructured data", 26, DN, True)
tx(s, Inches(1), Inches(3.2), Inches(5), Inches(2.5), "PRDs, BRDs, RFPs, meeting notes - enterprises drown in documents containing critical business intelligence. Extracting structure, KPIs, workflows, and risks is manual, slow, and error-prone.", 14, S)
pts = [("Days to weeks for manual analysis", R), ("No standardized business schema", O), ("Reports depend on individual consultants", O), ("Insights lost between projects", O)]
for i,(pt,clr) in enumerate(pts):
    y = 2.0 + i*0.7
    r(s, Inches(7.5), Inches(y), Inches(0.06), Inches(0.45), clr)
    tx(s, Inches(7.9), Inches(y), Inches(4.5), Inches(0.45), pt, 14, DT, True)
pn(s, 1)

# ==== S3: SOLUTION ====
s = blank(); r(s, Inches(0), Inches(0), SW, SH, PG); hdr(s, "The Solution", "02")
tx(s, Inches(1), Inches(1.5), Inches(11), Inches(0.8), "One input. Automatic intelligence.", 30, DN, True)
tx(s, Inches(1), Inches(2.0), Inches(10), Inches(0.5), "Input any business document and receive a complete, consulting-grade business system", 14, S)
steps = [("01", "Input", "PRD / BRD /\nMeeting Notes"), ("02", "Compile", "Business\nStructure"), ("03", "Analyze", "KPI System\n+ Workflows"), ("04", "Report", "PPT + HTML\nDashboard")]
for i,(num,stp,desc) in enumerate(steps):
    x = 1.5 + i*2.5
    r(s, Inches(x), Inches(3.2), Inches(2.2), Inches(2.5), W)
    r(s, Inches(x), Inches(3.2), Inches(2.2), Inches(0.06), TB)
    tx(s, Inches(x+0.2), Inches(3.4), Inches(1.8), Inches(0.5), num, 32, NC, True, "Arial")
    tx(s, Inches(x+0.2), Inches(3.9), Inches(1.8), Inches(0.4), stp, 16, DN, True)
    tx(s, Inches(x+0.2), Inches(4.4), Inches(1.8), Inches(1), desc, 11, S)
    if i < 3: tx(s, Inches(x+2.3), Inches(4.1), Inches(0.3), Inches(0.4), "\u2192", 20, TB, True, "Arial")
pn(s, 2)

# ==== S4: ARCHITECTURE ====
s = blank(); r(s, Inches(0), Inches(0), SW, SH, PG); hdr(s, "Technical Architecture", "03")
layers = [("Input Layer", "PRD / BRD / RFP / Meeting Notes / Documents", DN), ("AI Orchestration", "Intent Detection > Goal Routing > Agent Dispatch", TB), ("5-Agent Pipeline", "Compiler > Knowledge > Graph > Decision > Asset", NC), ("Output Layer", "PPT / HTML Dashboard / PDF / Word / JSON / XLSX", G)]
for i,(nm,dsc,clr) in enumerate(layers):
    y = 1.6 + i*1.15
    r(s, Inches(1), Inches(y), Inches(11), Inches(0.95), W)
    r(s, Inches(1), Inches(y), Inches(0.06), Inches(0.95), clr)
    tx(s, Inches(1.5), Inches(y+0.1), Inches(3), Inches(0.4), nm, 16, clr, True)
    tx(s, Inches(1.5), Inches(y+0.5), Inches(10), Inches(0.4), dsc, 12, S, False, "Arial")
kpi(s, Inches(1), Inches(6.3), Inches(2.5), Inches(0.8), "API Response", "<100ms", TB)
kpi(s, Inches(3.8), Inches(6.3), Inches(2.5), Inches(0.8), "Agent Pipeline", "~70ms", G)
kpi(s, Inches(6.6), Inches(6.3), Inches(2.5), Inches(0.8), "PPT Gen", "~40ms", NC)
kpi(s, Inches(9.4), Inches(6.3), Inches(2.5), Inches(0.8), "Tests Passed", "100%", G)
pn(s, 3)

# ==== S5: AGENT PIPELINE ====
s = blank(); r(s, Inches(0), Inches(0), SW, SH, PG); hdr(s, "Agent Pipeline", "04")
agents = [("Compiler Agent", "Parses documents, extracts business structure, modules, processes", "1.2ms", TB), ("Knowledge Agent", "Builds enterprise knowledge graph with versioning", "50ms", NC), ("Graph Agent", "Constructs business relationship graph (nodes + edges)", "0.3ms", G), ("Decision Agent", "Health scoring, bottleneck analysis, risk heatmaps", "0.1ms", O), ("Asset Agent", "Generates PPT, HTML Dashboard, PDF, Word, Excel, Bid Proposals", "125ms", R)]
for i,(nm,dsc,lat,clr) in enumerate(agents):
    y = 1.5 + i*1.0
    r(s, Inches(1), Inches(y), Inches(11), Inches(0.8), W)
    r(s, Inches(1), Inches(y), Inches(0.06), Inches(0.8), clr)
    tx(s, Inches(1.5), Inches(y+0.05), Inches(3.5), Inches(0.4), f"0{i+1}  {nm}", 16, DN, True)
    tx(s, Inches(1.5), Inches(y+0.4), Inches(8), Inches(0.3), dsc, 11, S)
    tx(s, Inches(10.5), Inches(y+0.15), Inches(1.5), Inches(0.5), lat, 14, NC, True, "Arial", PP_ALIGN.RIGHT)
pn(s, 4)

# ==== S6: KEY METRICS (Dark) ====
s = blank(); r(s, Inches(0), Inches(0), SW, SH, DN)
r(s, Inches(0), Inches(0), SW, Inches(0.05), NC)
tx(s, Inches(1), Inches(0.5), Inches(6), Inches(0.6), "Key Metrics", 32, W, True)
mets = [("12", "Intent Types", "40+ keywords"), ("5", "AI Agents", "~70ms pipeline"), ("6", "Industry Templates", "Content to Fintech"), ("34", "Tests", "100% pass rate"), ("39KB", "PPT Output", "7-slide deck"), ("74KB", "HTML Dashboard", "Interactive + Chart.js")]
for i,(num,lab,sub) in enumerate(mets):
    col=i%3; row=i//3
    x=1.5+col*3.5; y=2.0+row*2.3
    r(s, Inches(x), Inches(y), Inches(3), Inches(1.8), RGBColor(0x0A,0x1A,0x32))
    tx(s, Inches(x+0.3), Inches(y+0.2), Inches(2.4), Inches(0.7), num, 42, NC, True, "Arial")
    tx(s, Inches(x+0.3), Inches(y+1.0), Inches(2.4), Inches(0.3), lab, 14, W, True)
    tx(s, Inches(x+0.3), Inches(y+1.3), Inches(2.4), Inches(0.3), sub, 10, RGBColor(0xA0,0xAE,0xC0), False, "Arial")
pn(s, 5)

# ==== S7: COMPETITIVE ADVANTAGES ====
s = blank(); r(s, Inches(0), Inches(0), SW, SH, PG); hdr(s, "Competitive Advantages", "05")
advs = [("One-Click Simplicity", "Plain language input. Zero learning curve. System handles classification, compilation, analysis, and output generation automatically.", TB), ("Consulting-Grade Output", "PPT decks follow McKinsey-style narrative structure. Reports read like top-tier consulting deliverables, not machine-generated text.", NC), ("Modular Agent Architecture", "Five independent agents with fail-safe execution. Each agent debuggable, testable, and upgradeable independently.", G), ("Enterprise Knowledge Graph", "Cross-project knowledge reuse. Past analyses inform future projects. Version-controlled intelligence that compounds over time.", O)]
for i,(ttl,dsc,clr) in enumerate(advs):
    y = 1.4 + i*1.4
    r(s, Inches(1), Inches(y), Inches(11), Inches(1.1), W)
    r(s, Inches(1), Inches(y), Inches(0.06), Inches(1.1), clr)
    tx(s, Inches(1.5), Inches(y+0.1), Inches(10), Inches(0.35), ttl, 18, DN, True)
    tx(s, Inches(1.5), Inches(y+0.5), Inches(10), Inches(0.5), dsc, 11, S)
pn(s, 6)

# ==== S8: MARKET ====
s = blank(); r(s, Inches(0), Inches(0), SW, SH, PG); hdr(s, "Market Opportunity", "06")
tx(s, Inches(1), Inches(1.5), Inches(11), Inches(0.8), "The enterprise business intelligence market is growing rapidly", 24, DN, True)
segs = [("Enterprise Consultants", "Automate proposals. Reduce turnaround from weeks to minutes.", "60%"), ("Bid / RFP Teams", "Auto-extract RFP structure, generate competitive bid proposals.", "25%"), ("Operations Managers", "Convert SOPs into live dashboards with KPI monitoring and risk alerts.", "10%"), ("Product Managers", "Turn PRDs into structured specifications with auto-generated metrics.", "5%")]
for i,(seg,dsc,pct) in enumerate(segs):
    y = 2.8 + i*1.1
    r(s, Inches(1), Inches(y), Inches(10.5), Inches(0.85), W)
    r(s, Inches(1), Inches(y), Inches(0.06), Inches(0.85), TB)
    tx(s, Inches(1.5), Inches(y+0.05), Inches(3), Inches(0.35), seg, 15, DN, True)
    tx(s, Inches(1.5), Inches(y+0.4), Inches(7.5), Inches(0.35), dsc, 10, S)
    tx(s, Inches(10), Inches(y+0.1), Inches(1.5), Inches(0.6), pct, 28, TB, True, "Arial", PP_ALIGN.RIGHT)
pn(s, 7)

# ==== S9: ROADMAP ====
s = blank(); r(s, Inches(0), Inches(0), SW, SH, PG); hdr(s, "Product Roadmap", "07")
phases = [("Q3 2026 - MVP", "Core compiler, 5-agent pipeline, REST API, SaaS UI, PPT + HTML output", G), ("Q4 2026 - v1.0", "Multi-tenancy, auth, project mgmt, template library, knowledge search", TB), ("Q1 2027 - v2.0", "Collaborative editing, custom model fine-tuning, enterprise SSO", NC)]
for i,(ph,dsc,clr) in enumerate(phases):
    x = 1.5 + i*3.5
    r(s, Inches(x), Inches(3.5), Inches(3.2), Inches(0.04), clr)
    r(s, Inches(x+1.5), Inches(3.3), Inches(0.2), Inches(0.2), clr)
    r(s, Inches(x), Inches(2), Inches(3.2), Inches(1.2), W)
    r(s, Inches(x), Inches(2), Inches(3.2), Inches(0.05), clr)
    tx(s, Inches(x+0.2), Inches(2.1), Inches(2.8), Inches(0.4), ph, 16, clr, True)
    tx(s, Inches(x+0.2), Inches(2.5), Inches(2.8), Inches(0.7), dsc, 10, S)
pn(s, 8)

# ==== S10: TEAM ====
s = blank(); r(s, Inches(0), Inches(0), SW, SH, PG); hdr(s, "Team", "08")
tx(s, Inches(1), Inches(1.5), Inches(11), Inches(0.7), "Full-stack AI engineering team with enterprise SaaS experience", 22, DN, True)
roles = ["AI/ML Engineer", "Full-stack Developer", "Product Designer", "Business Analyst"]
for i,role in enumerate(roles):
    x = 1.5 + i*2.8
    r(s, Inches(x), Inches(3.0), Inches(2.4), Inches(2.8), W)
    r(s, Inches(x), Inches(3.0), Inches(2.4), Inches(0.05), TB)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.7), Inches(3.3), Inches(1), Inches(1))
    c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xE5,0xE7,0xEB); c.line.fill.background()
    tx(s, Inches(x+0.1), Inches(4.5), Inches(2.2), Inches(0.4), role, 13, DN, True, a=PP_ALIGN.CENTER)
    tx(s, Inches(x+0.1), Inches(5.0), Inches(2.2), Inches(0.3), "TBD", 10, S, False, "Arial", PP_ALIGN.CENTER)
pn(s, 9)

# ==== S11: INVESTMENT ====
s = blank(); r(s, Inches(0), Inches(0), SW, SH, DN)
tx(s, Inches(1), Inches(1.0), Inches(6), Inches(0.6), "Investment Opportunity", 32, W, True)
r(s, Inches(1), Inches(1.6), Inches(0.8), Inches(0.05), NC)
items = [("Engineering", "Core platform development, multi-tenancy, enterprise features", "45%", TB), ("Go-to-Market", "Sales, marketing, partnerships, content", "30%", NC), ("AI/ML R&D", "Model fine-tuning, accuracy, industry specialization", "15%", G), ("Operations", "Infrastructure, legal, admin", "10%", O)]
for i,(item,dsc,pct,clr) in enumerate(items):
    y = 2.5 + i*1.1
    w = 10.5*(float(pct[:-1])/100)
    r(s, Inches(1), Inches(y), Inches(w), Inches(0.8), clr)
    tx(s, Inches(1.3), Inches(y+0.05), Inches(3), Inches(0.35), item, 16, W, True)
    tx(s, Inches(1.3), Inches(y+0.4), Inches(8), Inches(0.3), dsc, 10, RGBColor(0xCC,0xCC,0xCC))
    tx(s, Inches(11.8), Inches(y+0.15), Inches(1), Inches(0.5), pct, 24, NC, True, "Arial", PP_ALIGN.RIGHT)
pn(s, 10)

# ==== S12: THANK YOU ====
s = blank(); r(s, Inches(0), Inches(0), SW, SH, DN)
for rad in [2.5, 2.8, 3.1]:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.4), Inches(2.2), Inches(rad), Inches(rad))
    c.fill.background(); c.line.color.rgb = TB if rad<3 else NC; c.line.width = Pt(0.5)
tx(s, Inches(1), Inches(3.5), Inches(11), Inches(1), "Thank You", 56, W, True, a=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(4.5), Inches(11), Inches(0.6), "BSC Studio \u00b7 Business System Compiler", 20, NC, False, "Arial", PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(5.3), Inches(11), Inches(0.4), "contact@bscstudio.com  |  bscstudio.com", 14, RGBColor(0xA0,0xAE,0xC0), False, "Arial", PP_ALIGN.CENTER)

# Save
out = r"C:\Users\34216\Documents\New project 3\bsc-backend\output\BSC_Studio_Investor_Pitch.pptx"
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f"DONE: {out}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {os.path.getsize(out):,} bytes")







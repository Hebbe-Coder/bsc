"""PRD Exporters - 多格式导出模块

基于PRD Document模型实现PDF/PPT/Word等多格式导出
"""
import os
import uuid
from typing import Optional


class PRDPDFExporter:
    """PRD PDF导出器"""
    
    def __init__(self):
        self._weasyprint_available = False
        self._pdfkit_available = False
        
        try:
            import weasyprint
            self._weasyprint_available = True
        except ImportError:
            pass
        
        try:
            import pdfkit
            self._pdfkit_available = True
        except ImportError:
            pass
    
    def export(self, document, output_path: Optional[str] = None) -> bytes:
        """
        导出PRD文档为PDF
        
        Args:
            document: PRDDocument对象
            output_path: 输出路径，若为None则返回字节流
        
        Returns:
            PDF字节流或输出路径
        """
        html_content = self._generate_html(document)
        
        if self._weasyprint_available:
            return self._export_with_weasyprint(html_content, output_path)
        elif self._pdfkit_available:
            return self._export_with_pdfkit(html_content, output_path)
        else:
            raise RuntimeError("需要安装weasyprint或pdfkit来生成PDF")
    
    def _generate_html(self, document) -> str:
        """生成PDF专用HTML"""
        return document.to_html()
    
    def _export_with_weasyprint(self, html_content: str, output_path: Optional[str]) -> bytes:
        """使用WeasyPrint生成PDF"""
        import weasyprint
        
        html = weasyprint.HTML(string=html_content)
        
        if output_path:
            html.write_pdf(output_path)
            return open(output_path, "rb").read()
        
        return html.write_pdf()
    
    def _export_with_pdfkit(self, html_content: str, output_path: Optional[str]) -> bytes:
        """使用pdfkit生成PDF"""
        import pdfkit
        
        options = {
            'encoding': 'UTF-8',
            'page-size': 'A4',
            'margin-top': '20mm',
            'margin-right': '20mm',
            'margin-bottom': '20mm',
            'margin-left': '20mm',
            'enable-local-file-access': None,
        }
        
        if output_path:
            pdfkit.from_string(html_content, output_path, options=options)
            return open(output_path, "rb").read()
        
        return pdfkit.from_string(html_content, False, options=options)


class PRDPPTExporter:
    """PRD PPT导出器"""
    
    def __init__(self):
        try:
            from pptx import Presentation
            self._pptx_available = True
        except ImportError:
            self._pptx_available = False
    
    def export(self, document, output_path: Optional[str] = None) -> str:
        """
        导出PRD文档为PPT
        
        Args:
            document: PRDDocument对象
            output_path: 输出路径，若为None则自动生成路径
        
        Returns:
            输出文件路径
        """
        if not self._pptx_available:
            raise RuntimeError("需要安装python-pptx来生成PPT")
        
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        self._add_cover_slide(prs, document)
        
        sections = document.sections
        for section in sections:
            self._add_section_slides(prs, section)
        
        if output_path is None:
            output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "output")
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, f"PRD_{document.title}_{uuid.uuid4().hex[:8]}.pptx")
        
        prs.save(output_path)
        return output_path
    
    def _add_cover_slide(self, prs, document):
        """添加封面幻灯片"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.333), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = document.title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 58, 138)
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = f"行业：{document.industry}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(107, 114, 128)
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = "产品需求文档"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(148, 163, 184)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_section_slides(self, prs, section, depth: int = 0):
        """递归添加章节幻灯片"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        left = Inches(1.0)
        top = Inches(1.0)
        width = Inches(11.333)
        height = Inches(5.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        title_p = tf.paragraphs[0]
        title_p.text = section.title
        title_p.font.size = Pt(20 if depth == 0 else 18)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(30, 58, 138)
        
        if section.content:
            content_p = tf.add_paragraph()
            content_p.text = self._clean_content(section.content)
            content_p.font.size = Pt(12)
            content_p.font.color.rgb = RGBColor(51, 65, 85)
            content_p.space_after = Pt(6)
        
        if section.children:
            for i, child in enumerate(section.children):
                child_p = tf.add_paragraph()
                child_p.text = f"{i+1}. {child.title}"
                child_p.font.size = Pt(14)
                child_p.font.bold = True
                child_p.font.color.rgb = RGBColor(59, 130, 246)
                child_p.space_before = Pt(12)
                
                if child.content:
                    child_content_p = tf.add_paragraph()
                    child_content_p.text = self._clean_content(child.content)[:300]
                    child_content_p.font.size = Pt(11)
                    child_content_p.font.color.rgb = RGBColor(75, 85, 99)
                    child_content_p.indent = Inches(0.3)
    
    def _clean_content(self, content: str) -> str:
        """清理内容，移除Markdown标记"""
        import re
        
        content = re.sub(r'#+\s', '', content)
        content = re.sub(r'\*\*([^*]+)\*\*', r'\1', content)
        content = re.sub(r'\*([^*]+)\*', r'\1', content)
        content = re.sub(r'`([^`]+)`', r'\1', content)
        content = re.sub(r'^\d+\.\s', '', content, flags=re.MULTILINE)
        content = re.sub(r'^\s*[-*+]\s', '', content, flags=re.MULTILINE)
        content = content.strip()
        
        return content


class PRDWordExporter:
    """PRD Word导出器"""
    
    def __init__(self):
        try:
            from docx import Document
            self._docx_available = True
        except ImportError:
            self._docx_available = False
    
    def export(self, document, output_path: Optional[str] = None) -> str:
        """
        导出PRD文档为Word
        
        Args:
            document: PRDDocument对象
            output_path: 输出路径，若为None则自动生成路径
        
        Returns:
            输出文件路径
        """
        if not self._docx_available:
            raise RuntimeError("需要安装python-docx来生成Word")
        
        from docx import Document
        from docx.shared import Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        
        doc = Document()
        
        style = doc.styles['Normal']
        font = style.font
        font.name = '微软雅黑'
        font.size = Pt(10.5)
        style.element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
        
        title_heading = doc.add_heading(document.title, level=0)
        title_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        info_paragraph = doc.add_paragraph()
        info_paragraph.add_run(f"行业：{document.industry} | 产品需求文档").italic = True
        info_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_page_break()
        
        self._add_sections_to_doc(doc, document.sections)
        
        if output_path is None:
            output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "output")
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, f"PRD_{document.title}_{uuid.uuid4().hex[:8]}.docx")
        
        doc.save(output_path)
        return output_path
    
    def _add_sections_to_doc(self, doc, sections, depth: int = 0):
        """递归添加章节到文档"""
        from docx.shared import Pt
        
        for section in sorted(sections, key=lambda s: s.order):
            heading_level = min(depth + 1, 9)
            heading = doc.add_heading(section.title, level=heading_level)
            
            if section.content:
                content = self._clean_content(section.content)
                paragraph = doc.add_paragraph(content)
                paragraph.paragraph_format.first_line_indent = Pt(21)
                paragraph.paragraph_format.line_spacing = 1.5
            
            if section.children:
                self._add_sections_to_doc(doc, section.children, depth + 1)
    
    def _clean_content(self, content: str) -> str:
        """清理内容，移除Markdown标记"""
        import re
        
        content = re.sub(r'#+\s', '', content)
        content = re.sub(r'\*\*([^*]+)\*\*', r'\1', content)
        content = re.sub(r'\*([^*]+)\*', r'\1', content)
        content = re.sub(r'`([^`]+)`', r'\1', content)
        content = re.sub(r'^\s*[-*+]\s', '', content, flags=re.MULTILINE)
        content = re.sub(r'^\d+\.\s', '', content, flags=re.MULTILINE)
        content = re.sub(r'^\s*>\s', '', content, flags=re.MULTILINE)
        content = re.sub(r'```[\s\S]*?```', '', content)
        content = content.strip()
        
        return content


def export_prd(document, format: str, output_path: Optional[str] = None):
    """
    统一导出接口
    
    Args:
        document: PRDDocument对象
        format: 导出格式（pdf/ppt/word/markdown/html）
        output_path: 输出路径
    
    Returns:
        PDF返回字节流，其他格式返回文件路径
    """
    format = format.lower()
    
    exporters = {
        'pdf': PRDPDFExporter,
        'ppt': PRDPPTExporter,
        'word': PRDWordExporter,
    }
    
    if format not in exporters:
        raise ValueError(f"不支持的导出格式：{format}")
    
    exporter = exporters[format]()
    return exporter.export(document, output_path)
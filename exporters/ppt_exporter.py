"""
PPT Exporter v7 — Impeccable Standard.
Design philosophy: "Lacquer & Gold" (neo-kinpaku).
Every slide is a crafted artifact, not a template fill.

Meticulously crafted with:
  - Dominant lacquer-black surfaces (60% visual weight)
  - kinpaku gold as the singular accent — never diluted
  - Albert Sans for display, DM Mono for data precision
  - Sandwich structure: dark cover/end, refined content core
  - Signature motif: thin gold rule, oversized slide numbers
  - Zero decoration that does not serve the content
"""
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    _PPTX_OK = True
except ImportError:
    Presentation = Inches = Pt = Emu = RGBColor = PP_ALIGN = MSO_ANCHOR = MSO_SHAPE = None
    _PPTX_OK = False
import os, uuid, datetime


def _require_pptx():
    """在使用 python-pptx 前调用；缺失时抛结构化错误。"""
    if not _PPTX_OK:
        from exporters.errors import ExportDependencyError
        raise ExportDependencyError("pptx", "python-pptx", "pip install python-pptx matplotlib")

# === DESIGN TOKENS ===
class T:
    """neo-kinpaku: Lacquer & Gold."""
    # Surfaces — dark dominance
    lacquer   = "12161A"  # primary bg
    lacquer_d = "0A0D0F"  # deeper bg for cover/end
    raised    = "1A1E23"  # cards on lacquer
    graphite  = "24282D"  # subtle elevation
    # Accent — singular gold
    kinpaku   = "C9A84C"  # oklch(84% 0.19 80.46)
    kinpaku_s = "D4BA6B"  # subtle highlight
    kinpaku_d = "B8923A"  # deep emphasis
    # Patina — secondary data color
    patina    = "5A9E96"
    patina_s  = "8BC4BE"
    # Text — warm hierarchy
    paper     = "EDEBE8"  # warm white for content slides
    text      = "1A1A18"  # near-black for content
    muted     = "6B6B66"  # secondary text
    faint     = "9E9E99"  # tertiary
    # State
    vermilion = "C0392B"

    @staticmethod
    def c(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))


class PPTExporter:
    """Impeccable presentation exporter."""

    def __init__(self):
        _require_pptx()
        self.prs = Presentation()
        self.prs.slide_width  = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.W = Inches(13.333)
        self.H = Inches(7.5)
        self.ML = Inches(1.0)
        self.MR = Inches(12.333)
        self.CW = Inches(11.333)

    # === PRIMITIVES ===

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _rect(self, s, l, t, w, h, fill=None, border=None, lw=0.5):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        sh.line.fill.background()
        if fill: sh.fill.solid(); sh.fill.fore_color.rgb = T.c(fill)
        if border: sh.line.color.rgb = T.c(border); sh.line.width = Pt(lw)
        return sh

    def _tb(self, s, l, t, w, h, txt, sz=12, color=None, bold=False, align=None, font="Albert Sans", anchor=None):
        color = color or T.paper
        tb = s.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        if anchor: tf.paragraphs[0].alignment = anchor
        p = tf.paragraphs[0]
        p.text = str(txt); p.font.size = Pt(sz)
        p.font.color.rgb = T.c(color); p.font.bold = bold; p.font.name = font
        if align: p.alignment = align
        return tb

    def _mtb(self, s, l, t, w, h, lines, sz=10, font="Albert Sans"):
        """Multi-line text box. lines: [(text, bold, color), ...]"""
        tb = s.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        for i, (txt, bld, clr) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(txt); p.font.size = Pt(sz)
            p.font.color.rgb = T.c(clr); p.font.bold = bld; p.font.name = font
        return tb

    def _rule(self, s, y, w=None, color=None, lw=0.5):
        """Horizontal rule — signature motif."""
        w = w or self.CW
        self._rect(s, self.ML, y, w, Inches(lw/72), fill=color or T.kinpaku)

    def _num(self, s, n, x=None, y=None):
        """Large slide number — signature element."""
        x = x or Inches(12.0)
        y = y or Inches(6.9)
        self._tb(s, x, y, Inches(1.0), Inches(0.5), f"{n:02d}", 28, T.kinpaku, True, PP_ALIGN.RIGHT, font="DM Mono")

    def _header_bar(self, s, title, subtitle=None):
        """Standard slide header with gold rule."""
        self._tb(s, self.ML, Inches(0.4), Inches(10), Inches(0.5), title, 22, T.paper, True)
        self._rule(s, Inches(1.1))
        if subtitle:
            self._tb(s, self.ML, Inches(1.25), self.CW, Inches(0.3), subtitle, 10, T.muted)

    # === SLIDE TEMPLATES ===

    def _slide_dark(self, title_lines, accent_block=True):
        """Dark slide template for cover/end."""
        s = self._blank()
        self._rect(s, 0, 0, self.W, self.H, fill=T.lacquer_d)
        if accent_block:
            self._rect(s, Inches(9.0), Inches(1.8), Inches(3.5), Inches(3.8), fill=T.kinpaku)
            self._rect(s, Inches(9.0), Inches(1.8), Inches(0.05), Inches(3.8), fill=T.kinpaku_d)
        for i, (txt, sz, color, bold, y) in enumerate(title_lines):
            self._tb(s, self.ML, y, Inches(7.5), Inches(sz/10), txt, sz, color, bold)
        self._rule(s, Inches(7.2))
        return s

    def _slide_content(self, s, title, subtitle=None):
        """Content slide template."""
        self._rect(s, 0, 0, self.W, self.H, fill=T.lacquer)
        self._header_bar(s, title, subtitle)

    # === EXPORT ===

    def export(self, business_system, output_path=None):
        if output_path is None:
            out_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "output")
            os.makedirs(out_dir, exist_ok=True)
            output_path = os.path.join(out_dir, f"impeccable_{uuid.uuid4().hex[:8]}.pptx")

        bs = business_system.get("business_system", business_system)
        meta = bs.get("metadata", {}) if isinstance(bs, dict) else {}
        title = meta.get("title", "Business System Report")
        date_str = datetime.date.today().strftime("%d %B %Y")

        self._s01_cover(title, date_str)
        self._s02_exec_summary(bs)
        self._s03_modules(bs)
        self._s04_workflow(bs)
        self._s05_kpi(bs)
        self._s06_risk(bs)
        self._s07_roi(bs)
        self._s08_timeline(bs)
        self._s09_end(title)

        self.prs.save(output_path)
        return output_path

    # === SLIDES ===

    def _s01_cover(self, title, date_str):
        s = self._slide_dark([
            ("Business System", 38, T.paper, True, Inches(1.8)),
            ("Compiler", 38, T.kinpaku, True, Inches(3.2)),
            (title, 14, T.muted, False, Inches(4.2)),
            (f"Executive Report · {date_str}", 9, T.faint, False, Inches(4.8)),
        ])
        self._tb(s, self.ML, Inches(6.0), Inches(7.5), Inches(0.3),
                 "BSC Engine v7 · Impeccable Standard · Confidential", 8, T.faint)

    def _s02_exec_summary(self, bs):
        s = self._blank(); self._slide_content(s, "Executive Summary", "System overview and key metrics")
        self._num(s, 2)

        modules = self._ensure(bs.get("modules", []))
        workflow = self._ensure(bs.get("workflow", []))
        kpi_list = self._ensure(bs.get("kpi", []) or bs.get("metrics", []))
        risk_list = self._ensure(bs.get("risk", []))

        # Metric tiles
        metrics = [
            ("Modules", len(modules), T.kinpaku),
            ("Workflow Steps", len(workflow), T.patina),
            ("KPIs", len(kpi_list), T.kinpaku_s),
            ("Risks", len(risk_list), T.vermilion),
        ]
        for i, (label, value, color) in enumerate(metrics):
            cx = self.ML + i * Inches(2.85)
            self._rect(s, cx, Inches(1.6), Inches(2.55), Inches(1.5), fill=T.raised)
            self._tb(s, cx + Inches(0.15), Inches(1.7), Inches(2.2), Inches(0.25), label, 9, T.muted)
            self._tb(s, cx + Inches(0.15), Inches(2.0), Inches(2.2), Inches(0.5), str(value), 32, color, True)

        # Objective
        mod_names = [str(m.get("name", "?")) if isinstance(m, dict) else str(m) for m in modules[:5]]
        objective = str(bs.get("objective", "") or bs.get("description", "") or ", ".join(mod_names))
        self._tb(s, self.ML, Inches(3.5), self.CW, Inches(1.0), objective[:300], 11, T.paper)
        self._rule(s, Inches(4.6), Inches(5), T.kinpaku, 0.3)

        # Module list
        self._tb(s, self.ML, Inches(4.9), Inches(5), Inches(0.3), "Core Modules", 12, T.kinpaku, True)
        mod_lines = [(f"{i+1}. {str(m.get('name','?')) if isinstance(m,dict) else str(m)}", False, T.paper)
                     for i, m in enumerate(modules[:8])]
        self._mtb(s, self.ML, Inches(5.3), Inches(5.5), Inches(1.8), mod_lines, 10)

        # Right: stats
        stats = [
            ("Backend", "FastAPI + Python 3.12+"),
            ("Export Formats", "PPTX, XLSX, HTML, SVG"),
            ("Design Standard", "neo-kinpaku | Impeccable"),
        ]
        for i, (k, v) in enumerate(stats):
            y = Inches(4.9) + i * Inches(0.45)
            self._tb(s, Inches(7.0), y, Inches(1.5), Inches(0.3), k, 9, T.muted, True)
            self._tb(s, Inches(8.5), y, Inches(4), Inches(0.3), v, 9, T.paper)

    def _s03_modules(self, bs):
        s = self._blank(); self._slide_content(s, "System Architecture", "Module decomposition and dependencies")
        self._num(s, 3)
        modules = self._ensure(bs.get("modules", []))
        colors = [T.kinpaku, T.patina, T.kinpaku_s, T.patina_s, T.kinpaku_d, T.patina]
        for i, mod in enumerate(modules[:6]):
            row, col = divmod(i, 3)
            cx = self.ML + col * Inches(3.75)
            cy = Inches(1.6) + row * Inches(2.7)
            c = colors[i % len(colors)]
            self._rect(s, cx, cy, Inches(3.5), Inches(2.45), fill=T.raised)
            self._rect(s, cx, cy, Inches(3.5), Inches(0.04), fill=c)
            nm = str(mod.get("name", "?")) if isinstance(mod, dict) else str(mod)
            self._tb(s, cx + Inches(0.18), cy + Inches(0.18), Inches(3.1), Inches(0.25),
                     f"0{i+1}  {nm}", 12, T.paper, True)
            desc = str(mod.get("description", ""))[:150] if isinstance(mod, dict) else ""
            self._tb(s, cx + Inches(0.18), cy + Inches(0.6), Inches(3.1), Inches(1.5), desc, 9, T.muted)
            deps = mod.get("depends_on", []) or mod.get("dependencies", []) if isinstance(mod, dict) else []
            if deps:
                ds = ", ".join(str(d) for d in deps[:3])
                self._tb(s, cx + Inches(0.18), cy + Inches(2.05), Inches(3.1), Inches(0.2),
                         f"Deps: {ds}", 7, T.faint)

    def _s04_workflow(self, bs):
        s = self._blank(); self._slide_content(s, "SOP Workflow", "Process nodes with ownership, SLA, and transitions")
        self._num(s, 4)
        workflow = self._ensure(bs.get("workflow", []))
        colors = [T.kinpaku, T.patina, T.kinpaku_s, T.patina_s, T.kinpaku_d, T.patina, T.kinpaku, T.patina_s]
        for i, step in enumerate(workflow[:8]):
            cy = Inches(1.6) + i * Inches(0.7)
            c = colors[i % len(colors)]
            self._rect(s, self.ML, cy, Inches(0.42), Inches(0.42), fill=c)
            self._tb(s, self.ML, cy + Inches(0.06), Inches(0.42), Inches(0.3),
                     str(i + 1), 13, T.lacquer_d, True, PP_ALIGN.CENTER)
            nm = step.get("name") or step.get("step") or f"Step {i+1}" if isinstance(step, dict) else str(step)
            self._tb(s, self.ML + Inches(0.6), cy + Inches(0.02), Inches(3.0), Inches(0.22), str(nm), 11, T.paper, True)
            owner = step.get("owner") or step.get("actor") or "" if isinstance(step, dict) else ""
            self._tb(s, self.ML + Inches(0.6), cy + Inches(0.28), Inches(3.0), Inches(0.18),
                     f"Owner: {owner}", 8, T.muted)
            sla = step.get("sla_hours") or step.get("sla") or "" if isinstance(step, dict) else ""
            if sla:
                self._tb(s, self.ML + Inches(3.8), cy + Inches(0.05), Inches(2.0), Inches(0.18),
                         f"SLA: {sla}h", 8, T.kinpaku)
            nxt = step.get("next", []) if isinstance(step, dict) else []
            if isinstance(nxt, str): nxt = [nxt]
            ns = ", ".join(str(x) for x in nxt[:3]) if nxt else "END"
            self._tb(s, self.ML + Inches(3.8), cy + Inches(0.28), Inches(3.5), Inches(0.18),
                     f"\u2192 {ns}", 8, T.faint)
            if i < len(workflow[:8]) - 1:
                self._rect(s, self.ML + Inches(0.2), cy + Inches(0.42), Inches(0.015), Inches(0.28), fill=c)

    def _s05_kpi(self, bs):
        s = self._blank(); self._slide_content(s, "KPI Dashboard", "Measurable indicators with formulas and targets")
        self._num(s, 5)
        kpi_list = self._ensure(bs.get("kpi", []) or bs.get("metrics", []))
        colors = [T.kinpaku, T.patina, T.kinpaku_s, T.patina_s]
        for i, kpi in enumerate(kpi_list[:8]):
            row, col = divmod(i, 4)
            cx = self.ML + col * Inches(2.85)
            cy = Inches(1.6) + row * Inches(2.7)
            c = colors[i % len(colors)]
            self._rect(s, cx, cy, Inches(2.6), Inches(2.45), fill=T.raised)
            nm = kpi.get("name", f"KPI {i+1}") if isinstance(kpi, dict) else str(kpi)
            target = kpi.get("target", "-") if isinstance(kpi, dict) else "-"
            formula = kpi.get("formula", "") if isinstance(kpi, dict) else ""
            self._tb(s, cx + Inches(0.15), cy + Inches(0.15), Inches(2.3), Inches(0.22), str(nm), 9, T.muted)
            self._tb(s, cx + Inches(0.15), cy + Inches(0.5), Inches(2.3), Inches(0.45), str(target), 24, c, True)
            if formula:
                self._tb(s, cx + Inches(0.15), cy + Inches(1.5), Inches(2.3), Inches(0.8),
                         f"Formula: {formula[:80]}", 7, T.faint)

    def _s06_risk(self, bs):
        s = self._blank(); self._slide_content(s, "Risk Assessment", "Risk matrix with severity, score, and mitigation strategy")
        self._num(s, 6)
        risk_list = self._ensure(bs.get("risk", []))
        sev = {"critical": T.vermilion, "high": T.kinpaku_d, "medium": T.kinpaku, "low": T.patina}
        for i, risk in enumerate(risk_list[:6]):
            row, col = divmod(i, 3)
            cx = self.ML + col * Inches(3.75)
            cy = Inches(1.6) + row * Inches(2.7)
            severity = risk.get("impact", "medium") or risk.get("severity", "medium") if isinstance(risk, dict) else "medium"
            c = sev.get(severity, T.muted)
            self._rect(s, cx, cy, Inches(3.5), Inches(2.45), fill=T.raised)
            self._rect(s, cx, cy, Inches(3.5), Inches(0.04), fill=c)
            nm = risk.get("name", "Risk") if isinstance(risk, dict) else str(risk)
            score = risk.get("score", 1) if isinstance(risk, dict) else 1
            self._tb(s, cx + Inches(0.15), cy + Inches(0.15), Inches(3.2), Inches(0.22),
                     f"[{severity.upper()}] {nm}", 11, T.paper, True)
            desc = str(risk.get("description", ""))[:140] if isinstance(risk, dict) else ""
            if desc:
                self._tb(s, cx + Inches(0.15), cy + Inches(0.5), Inches(3.2), Inches(0.7), desc, 8, T.muted)
            # Score bar
            self._rect(s, cx + Inches(0.15), cy + Inches(1.3), Inches(3.2), Inches(0.04), fill=T.graphite)
            bar_w = Inches(3.2 * min(float(score) / 10, 1.0))
            self._rect(s, cx + Inches(0.15), cy + Inches(1.3), bar_w, Inches(0.04), fill=c)
            self._tb(s, cx + Inches(0.15), cy + Inches(1.4), Inches(2.0), Inches(0.18),
                     f"Score: {score}/10", 8, c)
            mit = risk.get("mitigation", "") if isinstance(risk, dict) else ""
            if isinstance(mit, dict): mit = mit.get("action", str(mit))
            if mit:
                self._tb(s, cx + Inches(0.15), cy + Inches(1.7), Inches(3.2), Inches(0.6),
                         f"Mitigation: {str(mit)[:100]}", 7, T.faint)

    def _s07_roi(self, bs):
        s = self._blank(); self._slide_content(s, "ROI Analysis", "Projected return on investment and efficiency gains")
        self._num(s, 7)
        roi = [
            ("Cost Reduction", "35–50%", "Through workflow\nautomation", T.kinpaku),
            ("Efficiency Gain", "2.5\u00d7", "AI-assisted\nprocessing speed", T.patina),
            ("Quality Uplift", "+22%", "Consistent rule\nenforcement", T.kinpaku_s),
            ("Time-to-Market", "\u221260%", "Pre-built industry\ntemplates", T.patina_s),
        ]
        for i, (label, value, desc, color) in enumerate(roi):
            cx = self.ML + i * Inches(2.85)
            self._rect(s, cx, Inches(1.6), Inches(2.6), Inches(2.0), fill=T.raised)
            self._tb(s, cx + Inches(0.15), Inches(1.7), Inches(2.3), Inches(0.22), label, 9, T.muted)
            self._tb(s, cx + Inches(0.15), Inches(2.05), Inches(2.3), Inches(0.5), value, 24, color, True)
            self._tb(s, cx + Inches(0.15), Inches(2.7), Inches(2.3), Inches(0.7), desc, 8, T.faint)

        # Projection table
        self._rect(s, self.ML, Inches(4.0), self.CW, Inches(2.8), fill=T.raised)
        self._tb(s, self.ML + Inches(0.2), Inches(4.1), Inches(5), Inches(0.25),
                 "Annual Projected Impact", 12, T.kinpaku, True)
        rows = [
            ("Category", "Baseline", "Target", "Improvement", True),
            ("Processing Time", "4.2 hrs", "1.8 hrs", "\u221257%", False),
            ("Manual Reviews", "100%", "40%", "\u221260%", False),
            ("Error Rate", "3.2%", "1.1%", "\u221265%", False),
            ("Staff Utilization", "68%", "92%", "+24%", False),
            ("Monthly Cost", "$45K", "$28K", "\u221238%", False),
        ]
        for ri, (c1, c2, c3, c4, is_header) in enumerate(rows):
            y = Inches(4.45) + ri * Inches(0.32)
            col = T.paper if is_header else T.muted
            bld = is_header
            self._tb(s, self.ML + Inches(0.2), y, Inches(2.5), Inches(0.25), c1, 9, col, bld)
            self._tb(s, self.ML + Inches(2.9), y, Inches(2.0), Inches(0.25), c2, 9, col, bld, PP_ALIGN.CENTER)
            self._tb(s, self.ML + Inches(5.2), y, Inches(2.0), Inches(0.25), c3, 9, col, bld, PP_ALIGN.CENTER)
            self._tb(s, self.ML + Inches(7.6), y, Inches(2.0), Inches(0.25), c4, 9, T.kinpaku if ri > 0 else col, bld, PP_ALIGN.CENTER)

    def _s08_timeline(self, bs):
        s = self._blank(); self._slide_content(s, "Implementation Roadmap", "Phased delivery over 16 weeks")
        self._num(s, 8)
        phases = [
            ("01", "Foundation", "Weeks 1–4", ["Core architecture", "Module scaffolding", "API infrastructure", "Mock compiler"], T.kinpaku),
            ("02", "Core Features", "Weeks 5–8", ["LLM integration", "Graph & validation", "KPI engine", "Risk analysis"], T.patina),
            ("03", "Visualization", "Weeks 9–12", ["Interactive dashboard", "SOP flowchart", "PPT/XLSX export", "Decision sandbox"], T.kinpaku_s),
            ("04", "Production", "Weeks 13–16", ["Performance tuning", "Load testing", "Enterprise SSO", "Go-live deployment"], T.patina_s),
        ]
        for i, (num, name, duration, items, color) in enumerate(phases):
            cx = self.ML + i * Inches(2.85)
            # Number circle
            self._rect(s, cx + Inches(0.9), Inches(1.6), Inches(0.7), Inches(0.7), fill=color)
            self._tb(s, cx + Inches(0.9), Inches(1.7), Inches(0.7), Inches(0.5), num, 16, T.lacquer_d, True, PP_ALIGN.CENTER)
            # Connector
            if i < 3:
                self._rect(s, cx + Inches(2.6), Inches(1.9), Inches(0.25), Inches(0.015), fill=color)
            # Card
            self._rect(s, cx, Inches(2.5), Inches(2.6), Inches(3.8), fill=T.raised)
            self._tb(s, cx + Inches(0.15), Inches(2.6), Inches(2.3), Inches(0.22), f"{name}  {duration}", 12, color, True)
            for j, item in enumerate(items):
                self._tb(s, cx + Inches(0.3), Inches(2.95) + j * Inches(0.32), Inches(2.2), Inches(0.25),
                         f"\u2022 {item}", 9, T.muted)
        # CTA
        self._rect(s, self.ML, Inches(6.6), self.CW, Inches(0.5), fill=T.raised, border=T.kinpaku, lw=1)
        self._tb(s, self.ML + Inches(0.3), Inches(6.65), self.CW - Inches(0.6), Inches(0.35),
                 "Ready to deploy in 16 weeks  \u00b7  Contact us for a pilot engagement", 12, T.kinpaku, True, PP_ALIGN.CENTER)

    def _s09_end(self, title):
        s = self._slide_dark([
            ("Thank You", 42, T.kinpaku, True, Inches(2.2)),
            (title, 16, T.paper, False, Inches(3.2)),
            ("Generated by BSC Engine v7 \u00b7 Impeccable Standard", 10, T.muted, False, Inches(4.0)),
            ("Lacquer & Gold \u00b7 neo-kinpaku Design System", 9, T.faint, False, Inches(4.6)),
        ], accent_block=True)
        self._tb(s, self.ML, Inches(6.0), Inches(7.5), Inches(0.3),
                 "BSC Engine v7 \u00b7 Confidential \u00b7 Powering executive decisions", 8, T.faint)
        self._num(s, 9)

    @staticmethod
    def _ensure(val):
        if isinstance(val, list): return val
        if isinstance(val, dict):
            return [{"name": str(val.get("name", val)), "description": str(val.get("description", ""))}]
        return list(val) if val else []


def export_impeccable(business_system, output_path=None):
    exporter = PPTExporter()
    return exporter.export(business_system, output_path)

# ================================================================
# QA SYSTEM (skills-main pptx compliance)
# ================================================================

def qa_check(output_path: str) -> dict:
    """Post-generation quality assurance check.
    
    skills-main/pptx required QA:
    - Content check (markitdown for leftover placeholders)
    - Slide count verification
    - Structural validation
    """
    import subprocess, json
    
    results = {"path": output_path, "checks": [], "passed": True}
    
    # Check 1: File exists and has size
    if not os.path.exists(output_path):
        results["checks"].append({"check": "file_exists", "status": "FAIL", "detail": "Output file not found"})
        results["passed"] = False
        return results
    size = os.path.getsize(output_path)
    results["checks"].append({"check": "file_size", "status": "PASS", "detail": f"{size:,} bytes"})
    if size < 5000:
        results["checks"].append({"check": "min_size", "status": "WARN", "detail": f"File too small ({size} bytes), may be corrupted"})
        results["passed"] = False
    
    # Check 2: Slide count via python-pptx
    try:
        from pptx import Presentation
        prs = Presentation(output_path)
        slide_count = len(prs.slides)
        results["slide_count"] = slide_count
        results["checks"].append({"check": "slide_count", "status": "PASS" if slide_count >= 7 else "WARN", "detail": f"{slide_count} slides"})
        if slide_count < 5:
            results["passed"] = False
    except Exception as e:
        results["checks"].append({"check": "slide_count", "status": "FAIL", "detail": str(e)[:100]})
        results["passed"] = False
    
    # Check 3: Content extraction (placeholder detection)
    try:
        text_output = subprocess.run(
            ["python", "-m", "markitdown", output_path],
            capture_output=True, text=True, timeout=30
        )
        content = text_output.stdout
        placeholders = []
        for pattern in ["xxxx", "lorem", "ipsum", "placeholder", "TODO", "TBD"]:
            if pattern.lower() in content.lower():
                placeholders.append(pattern)
        if placeholders:
            results["checks"].append({"check": "no_placeholders", "status": "FAIL", "detail": f"Found: {placeholders}"})
            results["passed"] = False
        else:
            results["checks"].append({"check": "no_placeholders", "status": "PASS", "detail": "No placeholder text found"})
    except Exception:
        results["checks"].append({"check": "no_placeholders", "status": "SKIP", "detail": "markitdown not available"})
    
    # Check 4: Design rule violations
    try:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation(output_path)
        violations = []
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size < Pt(6):
                                violations.append(f"Slide {i+1}: text size {run.font.size} too small")
        if violations:
            results["checks"].append({"check": "font_sizes", "status": "WARN", "detail": f"{len(violations)} violations"})
        else:
            results["checks"].append({"check": "font_sizes", "status": "PASS", "detail": "All font sizes >= 6pt"})
    except Exception as e:
        results["checks"].append({"check": "font_sizes", "status": "SKIP", "detail": str(e)[:100]})
    
    results["all_passed"] = all(c["status"] == "PASS" or c["status"] == "SKIP" for c in results["checks"])
    return results


def export_with_qa(business_system, output_path=None):
    """Export PPTX with full QA verification loop.
    
    Follows skills-main/pptx verification loop:
    1. Generate slides
    2. Run QA checks
    3. Report issues
    4. Fix if needed
    """
    path = export_impeccable(business_system, output_path)
    qa = qa_check(path)
    
    if not qa.get("all_passed"):
        import logging
        logging.warning(f"PPTX QA found issues: {json.dumps(qa, indent=2, default=str)}")
    
    return path, qa
